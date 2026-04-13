"""在 LoRA微调实战技术分享_improved.pptx 中插入一页 '优化实录' slide"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import json, copy
from pathlib import Path
from lxml import etree

PPTX_PATH = Path(__file__).parent / "LoRA微调实战技术分享_improved.pptx"
REPORT = Path(__file__).parent.parent / "output" / "compare_report.json"

ACCENT = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x5F, 0x63, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x0D, 0x90, 0x4F)
RED = RGBColor(0xD9, 0x30, 0x25)
ORANGE = RGBColor(0xE3, 0x74, 0x00)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
BG_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)


def add_textbox(slide, left, top, width, height, text,
                font_size=18, bold=False, color=DARK,
                alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items,
                font_size=14, color=DARK, spacing=Pt(4)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(0.32 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = "Microsoft YaHei"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
    return table_shape


def move_slide_to_position(prs, target_index):
    """Move the last slide to a specific position (0-based index)."""
    ns = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = prs.part._element.find('.//p:sldIdLst', ns)
    children = list(sldIdLst)
    last_el = children[-1]
    sldIdLst.remove(last_el)
    if target_index >= len(children) - 1:
        sldIdLst.append(last_el)
    else:
        sldIdLst.insert(target_index, last_el)


def build():
    prs = Presentation(str(PPTX_PATH))

    with open(REPORT) as f:
        report = json.load(f)
    metrics = report["metrics"]
    m_base = metrics["No LoRA (基座)"]
    m_small = metrics["Small LoRA (30条)"]
    m_large = metrics["Large LoRA (200条)"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # BLANK

    # ---------- Title bar (match existing style) ----------
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.75))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = ACCENT
    title_bar.line.fill.background()

    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.5),
                "实战优化实录：数据量 × 训练策略的效果验证",
                font_size=22, bold=True, color=WHITE)

    # ---------- Left column: comparison table ----------
    add_textbox(slide, Inches(0.5), Inches(0.9), Inches(5), Inches(0.35),
                "三组配置 — 量化对比", font_size=13, bold=True, color=ACCENT)

    add_table(slide, Inches(0.5), Inches(1.25), Inches(5.8), [
        ["配置", "PPL ↓", "黑话命中 ↑", "数据量", "训练轮次", "Best Loss"],
        ["No LoRA (基座)", f"{m_base['ppl']:.1f}", f"{m_base['slang_hit']:.1%}", "—", "—", "—"],
        ["Small LoRA", f"{m_small['ppl']:.2f}", f"{m_small['slang_hit']:.1%}", "30条", "~8", "~1.08"],
        ["Large LoRA", f"{m_large['ppl']:.2f}", f"{m_large['slang_hit']:.1%}", "200条", "29", "0.065"],
    ], col_widths=[Inches(1.3), Inches(0.7), Inches(0.9), Inches(0.7), Inches(0.8), Inches(0.8)])

    # ---------- Right column: key finding card ----------
    add_rounded_rect(slide, Inches(6.7), Inches(0.9), Inches(5.8), Inches(1.7),
                     RGBColor(0xFC, 0xE8, 0xE6))

    add_textbox(slide, Inches(6.9), Inches(0.95), Inches(5.4), Inches(0.3),
                "⚠️ 核心发现", font_size=13, bold=True, color=RED)

    add_bullets(slide, Inches(6.9), Inches(1.3), Inches(5.4), Inches(1.2), [
        "6.7× 数据量 (30→200)，黑话命中率仅 +3.5%",
        "PPL 均 ≈1.0 → 过拟合信号 (模型在「背」而非「学」)",
        "根因：1.8B 小模型 + LoRA r=16 容量天花板",
        "优化 lr/dropout/delta 后训练更深 (29 epochs vs 5)",
    ], font_size=11, color=DARK, spacing=Pt(3))

    # ---------- Bottom section: two side-by-side panels ----------
    panel_top = Inches(3.0)

    # Left panel: optimization strategy
    add_rounded_rect(slide, Inches(0.5), panel_top, Inches(5.8), Inches(3.9), BG_LIGHT)

    add_textbox(slide, Inches(0.7), panel_top + Inches(0.08), Inches(5.4), Inches(0.3),
                "💡 训练策略优化", font_size=13, bold=True, color=ACCENT)

    add_table(slide, Inches(0.7), panel_top + Inches(0.45), Inches(5.4), [
        ["参数", "旧值", "新值", "为什么改"],
        ["learning_rate", "3e-4", "2e-4", "减速 → 渐进学习，避免暴跌式记忆"],
        ["lora_dropout", "0.1", "0.05", "200条数据充足 → 降低正则化"],
        ["MIN_DELTA", "0.01", "0.005", "更严格 early stop → 训练更多轮次"],
    ], col_widths=[Inches(1.0), Inches(0.6), Inches(0.6), Inches(3.2)])

    add_textbox(slide, Inches(0.7), panel_top + Inches(2.15), Inches(5.4), Inches(0.3),
                "训练效果: Loss 下降轨迹 (新 Large)", font_size=11, bold=True, color=ACCENT)

    add_table(slide, Inches(0.7), panel_top + Inches(2.45), Inches(5.4), [
        ["Epoch", "1", "5", "10", "15", "20", "25", "29"],
        ["Loss", "4.35", "1.88", "0.99", "0.29", "0.10", "0.066", "0.063"],
    ], col_widths=[Inches(0.7)] + [Inches(0.67)] * 7)

    add_textbox(slide, Inches(0.7), panel_top + Inches(3.25), Inches(5.4), Inches(0.5),
                "旧 Large ~5min / 5 epochs → 新 Large ~75min / 29 epochs，用时间换质量",
                font_size=10, color=GRAY)

    # Right panel: insights & takeaway
    add_rounded_rect(slide, Inches(6.7), panel_top, Inches(5.8), Inches(3.9), BG_ACCENT)

    add_textbox(slide, Inches(6.9), panel_top + Inches(0.08), Inches(5.4), Inches(0.3),
                "🧠 关键洞察", font_size=13, bold=True, color=ACCENT)

    insights = [
        "🏋️ 容量天花板 — 1.8B + r=16 仅 6.3M 可训参数，30条已接近上限",
        "📊 PPL ≈ 1.0 是过拟合而非好结果 — 需配合 LLM-as-Judge 评估",
        "🎲 生成随机性 (temp=0.7) 导致 3-5% 差距在噪声范围内",
        "📉 数据收益递减 — 6.7× 数据仅 +3.5% 命中率",
        "🗂️ 工程优化 — 移除冗余 tokenizer 文件，output 瘦身 31% (105→72MB)",
    ]
    add_bullets(slide, Inches(6.9), Inches(0.45) + panel_top, Inches(5.4), Inches(2.2),
                insights, font_size=11, color=DARK, spacing=Pt(6))

    # Takeaway bar at bottom of right panel
    add_rounded_rect(slide, Inches(7.0), panel_top + Inches(2.7), Inches(5.2), Inches(1.0), ACCENT)
    add_textbox(slide, Inches(7.2), panel_top + Inches(2.75), Inches(4.8), Inches(0.9),
                "结论: 数据量 ≠ 效果线性提升\n突破天花板 → 换更大模型 / 更优质数据 / 更强评测体系",
                font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ---------- Move new slide to position 24 (after Slide 24, 0-indexed = index 24) ----------
    move_slide_to_position(prs, 24)

    prs.save(str(PPTX_PATH))
    print(f"✅ 已在 Slide 25 位置插入「优化实录」页")
    print(f"   总页数: {len(prs.slides)}")
    print(f"   文件: {PPTX_PATH}")


if __name__ == "__main__":
    build()
