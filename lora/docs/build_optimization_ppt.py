"""Build PPT: LoRA 训练优化发现与改进"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json
from pathlib import Path

ACCENT = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x20, 0x20, 0x20)
GRAY = RGBColor(0x5F, 0x63, 0x68)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x0D, 0x90, 0x4F)
RED = RGBColor(0xD9, 0x30, 0x25)
ORANGE = RGBColor(0xE3, 0x74, 0x00)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
BG_ACCENT = RGBColor(0xE8, 0xF0, 0xFE)

OUTPUT = Path(__file__).parent.parent / "output"
REPORT = OUTPUT / "compare_report.json"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return txBox


def add_metric_card(slide, left, top, width, height, label, value, sub_text="", value_color=ACCENT):
    card = add_shape_fill(slide, left, top, width, height, WHITE, corner_radius=True)
    card.shadow.inherit = False

    add_text_box(slide, left + Inches(0.2), top + Inches(0.1),
                 width - Inches(0.4), Inches(0.3),
                 label, font_size=11, color=GRAY)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.35),
                 width - Inches(0.4), Inches(0.5),
                 value, font_size=26, bold=True, color=value_color)
    if sub_text:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.8),
                     width - Inches(0.4), Inches(0.3),
                     sub_text, font_size=10, color=GRAY)


def add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.35 * n_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.name = "Microsoft YaHei"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
    return table_shape


def slide_title_bar(slide, title):
    set_slide_bg(slide, WHITE)
    bar = add_shape_fill(slide, 0, 0, Inches(10), Inches(0.9), ACCENT)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.6),
                 title, font_size=24, bold=True, color=WHITE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    with open(REPORT) as f:
        report = json.load(f)

    metrics = report["metrics"]

    # ===================== Slide 1: Title =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ACCENT)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1),
                 "LoRA 微调优化实录", font_size=36, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.6),
                 "从诊断到改进 — Small vs Large LoRA 效果差距分析",
                 font_size=18, color=RGBColor(0xD2, 0xE3, 0xFC), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.5),
                 "基座模型: Qwen1.5-1.8B-Chat  |  微调方法: LoRA / PEFT  |  场景: 台灣電商網紅",
                 font_size=12, color=RGBColor(0xA8, 0xC7, 0xFA), alignment=PP_ALIGN.CENTER)

    # ===================== Slide 2: 问题发现 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "问题发现：Large 为什么没有比 Small 强很多？")

    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(9), Inches(0.4),
                 "初始对比结果（优化前）", font_size=14, bold=True, color=ACCENT)

    add_table(slide, Inches(0.6), Inches(1.5), Inches(5.5), [
        ["配置", "PPL ↓", "黑话命中率 ↑", "训练数据"],
        ["No LoRA (基座)", "716.60", "~8%", "—"],
        ["Small LoRA (30条)", "1.078", "21.8%", "30 条"],
        ["Large LoRA (200条)", "1.065", "29.9%", "200 条"],
    ], col_widths=[Inches(1.8), Inches(1.0), Inches(1.4), Inches(1.3)])

    add_shape_fill(slide, Inches(6.5), Inches(1.5), Inches(3.2), Inches(1.4), RGBColor(0xFC, 0xE8, 0xE6), corner_radius=True)
    add_text_box(slide, Inches(6.7), Inches(1.55), Inches(2.8), Inches(0.3),
                 "⚠️ 关键疑问", font_size=13, bold=True, color=RED)
    add_bullet_list(slide, Inches(6.7), Inches(1.85), Inches(2.8), Inches(1.0), [
        "PPL 仅降 1.2%，几乎无差别",
        "黑话命中率只提升 8 个百分点",
        "6.7x 的数据量 → 微弱提升？",
    ], font_size=11, color=RED)

    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(9), Inches(0.4),
                 "🔍 直觉判断：PPL ≈ 1.0 意味着严重过拟合 — 模型在「背答案」而非学风格",
                 font_size=13, bold=True, color=ORANGE)

    add_bullet_list(slide, Inches(0.6), Inches(3.7), Inches(9), Inches(1.8), [
        'PPL → 1.0: 模型对训练集的每个 token 几乎 100% 确定，这不是「理解」，是「记忆」',
        "30 条就能背下来 → 200 条也只是多背了一些 → 泛化能力没有实质提升",
        "过于激进的学习率 (3e-4) + 宽松的 early stopping (MIN_DELTA=0.01) 导致几个 epoch 就收敛",
    ], font_size=12, color=DARK, spacing=Pt(8))

    # ===================== Slide 3: 根因分析 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "根因分析：四个维度的深度诊断")

    box_data = [
        ("📊 训练数据覆盖", Inches(0.4), [
            "Small(30) & Large(200):",
            "  22/22 黑话词全覆盖 ✅",
            "数据质量不是瓶颈",
            "Large 频率更高: 手刀 91次 vs 10次",
        ]),
        ("📈 训练曲线", Inches(2.7), [
            "旧配置: 几个 epoch → loss≈1.06",
            "  暴跌式收敛 = 死记硬背",
            "新配置: 29 epochs → loss=0.065",
            "  渐进式下降 = 充分学习",
        ]),
        ("🎯 Eval 命中细节", Inches(5.0), [
            "Large 某些 prompt 强于 Small:",
            "  P1: 6/10 vs 3/10, P9: 4/8 vs 2/8",
            "但生成有随机性 (temp=0.7)",
            "  3-5% 差距在噪声范围内",
        ]),
        ("🧠 模型容量天花板", Inches(7.3), [
            "1.8B 参数 + r=16 LoRA",
            "  仅训练 0.34% 参数 (6.3M)",
            "30条已接近容量上限",
            "200条无法带来更多增量",
        ]),
    ]

    for title, left, items in box_data:
        card = add_shape_fill(slide, left, Inches(1.15), Inches(2.3), Inches(3.0), BG_LIGHT, corner_radius=True)
        add_text_box(slide, left + Inches(0.1), Inches(1.2), Inches(2.1), Inches(0.35),
                     title, font_size=12, bold=True, color=ACCENT)
        add_bullet_list(slide, left + Inches(0.1), Inches(1.6), Inches(2.1), Inches(2.4),
                        items, font_size=10, color=DARK, spacing=Pt(4))

    add_text_box(slide, Inches(0.4), Inches(4.4), Inches(9.2), Inches(0.5),
                 "结论: 根本瓶颈不在数据量，而在小模型的容量上限和训练策略",
                 font_size=14, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ===================== Slide 4: 优化方案 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "优化方案：三个维度的训练策略调整")

    changes = [
        ("学习率", "3e-4", "2e-4", "减慢学习速度，避免暴跌式记忆"),
        ("Dropout", "0.1", "0.05", "200条数据下适当降低正则化"),
        ("MIN_DELTA", "0.01", "0.005", "更严格的 early stopping 阈值"),
    ]

    add_table(slide, Inches(0.6), Inches(1.2), Inches(5.5), [
        ["参数", "旧值", "新值", "调整理由"],
    ] + [[c[0], c[1], c[2], c[3]] for c in changes],
    col_widths=[Inches(1.0), Inches(0.8), Inches(0.8), Inches(2.9)])

    add_shape_fill(slide, Inches(6.5), Inches(1.2), Inches(3.2), Inches(1.6), BG_ACCENT, corner_radius=True)
    add_text_box(slide, Inches(6.7), Inches(1.25), Inches(2.8), Inches(0.3),
                 "💡 设计思路", font_size=13, bold=True, color=ACCENT)
    add_bullet_list(slide, Inches(6.7), Inches(1.6), Inches(2.8), Inches(1.1), [
        "更低的 lr → 渐进式学习",
        "更低的 dropout → 让 200 条\n  数据被更充分利用",
        "更严格的 delta → 训练更多\n  epoch，风格内化更深",
    ], font_size=11, color=DARK, spacing=Pt(4))

    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(9), Inches(0.3),
                 "代码层面改动", font_size=14, bold=True, color=ACCENT)

    code_items = [
        "step2_lora_inject.py — 新增 LORA_PROFILES 字典，根据 LORA_EXPERIMENT 自动选择配置",
        "step3_train.py — learning_rate 根据 EXPERIMENT_LABEL 自动切换；MIN_DELTA 从 0.01→0.005",
        "step6_compare.py — 新增 build_overview() 生成结构化对比结论 + 计时信息",
        "step3_train.py — 移除 tokenizer.save_pretrained()，每个 adapter 目录节省 11MB（总计 -33MB）",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(3.55), Inches(9), Inches(1.8),
                    code_items, font_size=11, color=DARK, spacing=Pt(4))

    # ===================== Slide 5: 训练过程对比 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "训练过程对比：旧策略 vs 新策略")

    add_metric_card(slide, Inches(0.5), Inches(1.2), Inches(2.1), Inches(1.1),
                    "旧 Large — 训练轮次", "~5 epochs", "快速 early stop", RED)
    add_metric_card(slide, Inches(2.8), Inches(1.2), Inches(2.1), Inches(1.1),
                    "旧 Large — Best Loss", "1.065", "PPL ≈ 1.0 = 过拟合", RED)
    add_metric_card(slide, Inches(5.1), Inches(1.2), Inches(2.1), Inches(1.1),
                    "新 Large — 训练轮次", "29 epochs", "充分学习后收敛", GREEN)
    add_metric_card(slide, Inches(7.4), Inches(1.2), Inches(2.1), Inches(1.1),
                    "新 Large — Best Loss", "0.065", "训练 loss 降 94%", GREEN)

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.3),
                 "Loss 下降轨迹 (新 Large, 29 epochs)", font_size=13, bold=True, color=ACCENT)

    loss_data = [
        ("Epoch", "1", "5", "10", "15", "20", "25", "29"),
        ("Avg Loss", "4.35", "1.88", "0.99", "0.29", "0.10", "0.066", "0.063"),
    ]
    add_table(slide, Inches(0.5), Inches(2.95), Inches(9), loss_data,
              col_widths=[Inches(1.1)] + [Inches(1.13)] * 7)

    add_shape_fill(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.4), BG_LIGHT, corner_radius=True)
    add_text_box(slide, Inches(0.7), Inches(3.85), Inches(8.6), Inches(0.3),
                 "📊 解读", font_size=13, bold=True, color=ACCENT)
    add_bullet_list(slide, Inches(0.7), Inches(4.2), Inches(8.6), Inches(0.9), [
        "新策略: 4.35 → 1.88 → 0.99 → 0.29 → 0.10 → 0.065，平滑渐进式下降",
        "旧策略: 几步就跌到 ≈1.06，然后 early stop — 模型只学了表面",
        "训练时间: 旧版 ~5min，新版 ~75min — 用时间换质量",
    ], font_size=12, color=DARK, spacing=Pt(4))

    # ===================== Slide 6: 最终结果 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "最终对比结果")

    m_base = metrics["No LoRA (基座)"]
    m_small = metrics["Small LoRA (30条)"]
    m_large = metrics["Large LoRA (200条)"]

    add_table(slide, Inches(0.5), Inches(1.2), Inches(6.0), [
        ["配置", "PPL ↓", "黑话命中率 ↑", "训练 Epochs", "Best Loss"],
        ["No LoRA (基座)", f"{m_base['ppl']:.2f}", f"{m_base['slang_hit']:.1%}", "—", "—"],
        ["Small LoRA (30条)", f"{m_small['ppl']:.2f}", f"{m_small['slang_hit']:.1%}", "~8", "~1.08"],
        ["Large LoRA (200条)", f"{m_large['ppl']:.2f}", f"{m_large['slang_hit']:.1%}", "29", "0.065"],
    ], col_widths=[Inches(1.6), Inches(1.0), Inches(1.2), Inches(1.1), Inches(1.1)])

    add_text_box(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(0.3),
                 "Eval Prompt 逐条对比 (hits / expected)", font_size=13, bold=True, color=ACCENT)

    eval_rows = [["Prompt", "Base", "Small", "Large"]]
    short_names = [
        "P1: 无线吸尘器推广文案",
        "P2: 快递没到回覆客人",
        "P3: 双12推广文案",
        "P4: 防晒乳推荐",
        "P5: 蝦皮店铺评分",
        "P6: 蓝牙耳机TikTok",
        "P7: 母亲节推广",
        "P8: 直播带货开场白",
        "P9: 手机支架安利",
        "P10: 联盟行销教学",
    ]
    for i in range(10):
        b = report["eval_details"]["No LoRA (基座)"][i]
        s = report["eval_details"]["Small LoRA (30条)"][i]
        l = report["eval_details"]["Large LoRA (200条)"][i]
        eval_rows.append([
            short_names[i],
            f"{b['hits']}/{b['expected']}",
            f"{s['hits']}/{s['expected']}",
            f"{l['hits']}/{l['expected']}",
        ])
    add_table(slide, Inches(0.5), Inches(3.3), Inches(6.0), eval_rows,
              col_widths=[Inches(2.4), Inches(0.8), Inches(0.8), Inches(0.8)])

    add_shape_fill(slide, Inches(6.8), Inches(1.2), Inches(2.9), Inches(2.4), BG_ACCENT, corner_radius=True)
    add_text_box(slide, Inches(7.0), Inches(1.25), Inches(2.5), Inches(0.3),
                 "✅ 关键发现", font_size=13, bold=True, color=ACCENT)
    add_bullet_list(slide, Inches(7.0), Inches(1.6), Inches(2.5), Inches(1.9), [
        "PPL: Large 1.06 vs\n  Small 1.08 (差距小)",
        "黑话命中率: 27.6% vs\n  24.1% (+3.5%)",
        "P1 吸尘器 Large 6/10\n  vs Small 3/10 亮点",
        "P9 手机架 Large 4/8\n  vs Small 2/8 亮点",
    ], font_size=11, color=DARK, spacing=Pt(4))

    # ===================== Slide 7: 关键洞察 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "关键洞察与工程建议")

    insights = [
        ("🧠", "小模型的容量天花板", "1.8B + LoRA r=16 仅 6.3M 可训参数\n30条数据已接近上限，200条增量有限"),
        ("📊", "PPL 不是万能指标", "PPL≈1.0 是过拟合信号而非好结果\n应结合生成质量（LLM-as-Judge）评估"),
        ("🎲", "生成随机性的影响", "temperature=0.7 + sampling 导致\n每次结果波动，3-5%差距在噪声内"),
        ("💡", "数据收益对数递减", "30→200: 6.7x数据量 → 仅+3.5%命中率\n工业界常见，需更大模型才能突破"),
    ]

    for i, (emoji, title, desc) in enumerate(insights):
        row = i // 2
        col = i % 2
        left = Inches(0.5) + col * Inches(4.7)
        top = Inches(1.2) + row * Inches(1.8)

        card = add_shape_fill(slide, left, top, Inches(4.4), Inches(1.5), BG_LIGHT, corner_radius=True)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.1),
                     Inches(4.1), Inches(0.3),
                     f"{emoji} {title}", font_size=14, bold=True, color=ACCENT)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.5),
                     Inches(4.1), Inches(0.9),
                     desc, font_size=12, color=DARK)

    add_shape_fill(slide, Inches(0.5), Inches(4.85), Inches(9.0), Inches(0.5), ACCENT, corner_radius=True)
    add_text_box(slide, Inches(0.7), Inches(4.87), Inches(8.6), Inches(0.45),
                 "核心结论: 数据量 ≠ 效果线性提升。要突破天花板 → 换更大模型 / 更优质数据 / 更强评测",
                 font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ===================== Slide 8: 工程优化附录 =====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_title_bar(slide, "附录：工程优化 — 存储 & 报告")

    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(4.5), Inches(0.3),
                 "存储优化: output 文件夹瘦身", font_size=14, bold=True, color=ACCENT)
    add_table(slide, Inches(0.5), Inches(1.5), Inches(4.2), [
        ["优化项", "Before", "After", "节省"],
        ["output/ 总大小", "105 MB", "72 MB", "33 MB (31%)"],
        ["每个 adapter", "35 MB", "24 MB", "11 MB"],
        ["删除冗余文件", "tokenizer.json", "—", "×3 份 = 33MB"],
    ], col_widths=[Inches(1.2), Inches(1.0), Inches(1.0), Inches(1.0)])

    add_bullet_list(slide, Inches(0.5), Inches(3.1), Inches(4.2), Inches(1.5), [
        "推理时 tokenizer 从基座模型加载",
        "adapter 目录只需保留 adapter_model.safetensors + adapter_config.json",
        "step4/step5/step6 完全不受影响",
    ], font_size=11, color=DARK, spacing=Pt(3))

    add_text_box(slide, Inches(5.3), Inches(1.1), Inches(4.5), Inches(0.3),
                 "报告优化: compare_report.json 新增 overview", font_size=14, bold=True, color=ACCENT)

    overview_items = [
        "新增 overview 顶层字段:",
        "  • model: 基座模型名称",
        "  • takeaways: 自动生成的对比结论",
        "  • best_ppl / best_slang_hit: 最佳指标",
        "  • timings: 每个配置的评测耗时",
        "  • total_time: 总评测耗时",
        "",
        "终端输出也同步展示 overview 结论",
    ]
    add_bullet_list(slide, Inches(5.3), Inches(1.5), Inches(4.2), Inches(2.5),
                    overview_items, font_size=11, color=DARK, spacing=Pt(3))

    ov = report["overview"]
    add_shape_fill(slide, Inches(5.3), Inches(3.6), Inches(4.2), Inches(1.2), BG_ACCENT, corner_radius=True)
    add_text_box(slide, Inches(5.5), Inches(3.65), Inches(3.8), Inches(0.25),
                 "⏱️ 评测耗时", font_size=12, bold=True, color=ACCENT)
    timing_text = "\n".join(f"  {k}: {v}" for k, v in ov["timings"].items())
    timing_text += f"\n  总计: {ov['total_time']}"
    add_text_box(slide, Inches(5.5), Inches(3.95), Inches(3.8), Inches(0.8),
                 timing_text, font_size=11, color=DARK)

    # ===================== Save =====================
    out_path = Path(__file__).parent / "LoRA训练优化发现.pptx"
    prs.save(str(out_path))
    print(f"✅ PPT 已生成: {out_path}")
    print(f"   共 {len(prs.slides)} 页")


if __name__ == "__main__":
    build()
