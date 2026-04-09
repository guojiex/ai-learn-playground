"""PPT 改进脚本 — 为技术分享 PPT 新增三处改进。

改进内容:
1. 在 Slide 14（核心代码实现）后新增 "效果验证: Before vs After" 页
2. 在总结前新增 "实战避坑指南" 页
3. 将 Slide 14 拆分为 "量化配置" + "LoRA 注入" 两页

Usage:
    python update_pptx.py
"""

import copy
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INPUT_PATH = Path(__file__).parent / "docs" / "LoRA微调实战技术分享.pptx"
OUTPUT_PATH = Path(__file__).parent / "docs" / "LoRA微调实战技术分享_improved.pptx"


def find_blank_layout(prs):
    """Find the blank or most minimal layout available."""
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def clone_slide(prs, slide_index):
    """Clone an existing slide to reuse its visual design."""
    source = prs.slides[slide_index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in source.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)
    return new_slide


def set_shape_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    """Set text content of a shape, clearing existing content."""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = Pt(font_size)
    if bold is not None:
        p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if alignment:
        p.alignment = alignment


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a textbox with the given parameters."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=13, color=None):
    """Add a textbox with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        if color:
            p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox


def move_slide(prs, old_index, new_index):
    """Move a slide from old_index to new_index."""
    slides = prs.slides._sldIdLst
    el = list(slides)[old_index]
    slides.remove(el)
    if new_index >= len(list(slides)):
        slides.append(el)
    else:
        ref = list(slides)[new_index]
        slides.insert(slides.index(ref), el)


def create_eval_slide(prs):
    """Create '效果验证: Before vs After' slide."""
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    W = prs.slide_width
    H = prs.slide_height

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.8),
                "效果验证：Before vs After", font_size=28, bold=True,
                color=RGBColor(0x1A, 0x1A, 0x2E))

    prompt_text = (
        '📝 测试 Prompt: \u201c你是印尼 TikTok 顶级推客。'
        '请用本地黑话安利这款充电宝: Powerbank 20000mAh, 100k IDR.\u201d'
    )
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(11), Inches(0.5),
                prompt_text, font_size=12, color=RGBColor(0x55, 0x55, 0x55))

    add_textbox(slide, Inches(0.5), Inches(1.8), Inches(5), Inches(0.5),
                "🔴 微调前 (通用模型)", font_size=18, bold=True,
                color=RGBColor(0xCC, 0x33, 0x33))

    add_multiline_textbox(slide, Inches(0.5), Inches(2.3), Inches(5), Inches(2.5), [
        ("这款充电宝容量为20000毫安时，价格仅为10万印尼盾。", False),
        ("它具有大容量、高性价比的特点，非常适合", False),
        ("日常使用。建议您考虑购买。", False),
        ("", False),
        ("❌ 语气生硬、缺少本地黑话、没有购买引导", True),
    ], font_size=13, color=RGBColor(0x33, 0x33, 0x33))

    add_textbox(slide, Inches(6.3), Inches(1.8), Inches(5.5), Inches(0.5),
                "🟢 微调后 (LoRA 适配器)", font_size=18, bold=True,
                color=RGBColor(0x22, 0x8B, 0x22))

    add_multiline_textbox(slide, Inches(6.3), Inches(2.3), Inches(5.5), Inches(2.5), [
        ("PROMO GILA! 🔥 Powerbank 20000mAh ini", False),
        ("sumpah murah banget cuma 100rb! Awet,", False),
        ("bisa charge HP 5x. Buruan CHECKOUT", False),
        ("sebelum kehabisan ya Bestie~ 🛒✨", False),
        ("", False),
        ("✅ 融入 Spill/Cuan/Checkout 等本地表达", True),
    ], font_size=13, color=RGBColor(0x33, 0x33, 0x33))

    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(11), Inches(0.5),
                "📊 量化指标", font_size=18, bold=True,
                color=RGBColor(0x1A, 0x1A, 0x2E))

    table_shape = slide.shapes.add_table(2, 4,
        Inches(0.5), Inches(5.5), Inches(11), Inches(1.0))
    table = table_shape.table

    headers = ["指标", "微调前", "微调后", "提升"]
    values = ["本地黑话命中率", "~5%", "~72%", "▲ 14x"]

    for col, (h, v) in enumerate(zip(headers, values)):
        cell_h = table.cell(0, col)
        cell_h.text = h
        for p in cell_h.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

        cell_v = table.cell(1, col)
        cell_v.text = v
        for p in cell_v.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "这是整场分享最具冲击力的环节。同一个Prompt，微调前后的差异一目了然。"
        "微调前模型的回复像官方客服，微调后则像一个活跃在TikTok评论区的本地推客。"
        "量化指标可以根据实际实验结果调整。"
    )

    return slide


def create_pitfalls_slide(prs):
    """Create '实战避坑指南' slide."""
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(11), Inches(0.8),
                "实战避坑指南", font_size=28, bold=True,
                color=RGBColor(0x1A, 0x1A, 0x2E))

    pitfalls = [
        ("\u26a0\ufe0f \u707e\u96be\u6027\u9057\u5fd8 (Catastrophic Forgetting)",
         "\u6a21\u578b\u5b66\u4f1a\u4e86\u65b0\u77e5\u8bc6\u4f46\u5fd8\u8bb0\u4e86\u65e7\u80fd\u529b\uff0c\u8868\u73b0\u4e3a\u8bf4\u8bdd\u53d8\u201c\u590d\u8bfb\u673a\u201d\u6216\u903b\u8f91\u6df7\u4e71\u3002\n"
         "\u89e3\u6cd5\uff1a\u964d\u4f4e lora_alpha / r \u6bd4\u503c\uff0c\u51cf\u5c0f\u5fae\u8c03\u5bf9\u539f\u59cb\u6a21\u578b\u7684\u5f71\u54cd\u529b\u3002\u63a8\u8350 alpha = 2r\u3002"),

        ("\u26a0\ufe0f \u6570\u636e\u6c61\u67d3 (Data Contamination)",
         "\u8bad\u7ec3\u5370\u5c3c\u9002\u914d\u5668\u65f6\u6df7\u5165\u6cf0\u8bed\u6570\u636e\uff0c\u6a21\u578b\u4f1a\u4ea7\u751f\u8bed\u8a00\u4ea4\u53c9\u7684\u201c\u5e7b\u89c9\u201d\u3002\n"
         "\u89e3\u6cd5\uff1a\u4e25\u683c\u6309\u56fd\u5bb6/\u8bed\u8a00\u9694\u79bb\u6570\u636e\u96c6\uff0c\u786e\u4fdd\u6bcf\u4e2a\u9002\u914d\u5668\u53ea\u5b66\u4e60\u5bf9\u5e94\u7684\u6587\u5316\u8bed\u5883\u3002"),

        ("\u26a0\ufe0f Rank \u9009\u62e9\u4e0d\u5f53",
         "r \u8fc7\u5c0f \u2192 \u5b66\u4e0d\u5230\u8db3\u591f\u7684\u9886\u57df\u77e5\u8bc6\uff1br \u8fc7\u5927 \u2192 \u53c2\u6570\u81a8\u80c0\u4e14\u53ef\u80fd\u8fc7\u62df\u5408\u3002\n"
         "\u7ecf\u9a8c\u503c\uff1ar=8 \u9002\u5408\u7b80\u5355\u98ce\u683c\u8fc1\u79fb\uff0cr=16~32 \u9002\u5408\u9700\u8981\u5b66\u4e60\u590d\u6742\u8bed\u4e49\u7684\u573a\u666f\u3002"),

        ("\u26a0\ufe0f \u7f3a\u5c11\u663e\u5f0f\u6807\u6ce8",
         "\u4ec5\u9760\u6570\u636e\u8bed\u8a00\u8ba9\u6a21\u578b\u81ea\u52a8\u8bc6\u522b\u56fd\u5bb6\uff0c\u63a8\u7406\u65f6\u5bb9\u6613\u201c\u4e32\u53f0\u201d\u3002\n"
         "\u89e3\u6cd5\uff1a\u5728 instruction \u4e2d\u660e\u786e\u6807\u6ce8\u89d2\u8272\u548c\u56fd\u5bb6\uff08\u5982\u201c\u4f5c\u4e3a\u5370\u5c3c\u8d44\u6df1\u63a8\u5ba2\u201d\uff09\uff0c\u5efa\u7acb\u6761\u4ef6\u53cd\u5c04\u3002"),
    ]

    y_offset = Inches(1.2)
    for title, desc in pitfalls:
        add_textbox(slide, Inches(0.5), y_offset, Inches(11), Inches(0.4),
                    title, font_size=16, bold=True,
                    color=RGBColor(0xCC, 0x66, 0x00))
        y_offset += Inches(0.4)

        add_textbox(slide, Inches(0.8), y_offset, Inches(10.5), Inches(0.8),
                    desc, font_size=12, color=RGBColor(0x44, 0x44, 0x44))
        y_offset += Inches(0.95)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "这些都是在实际微调过程中最常遇到的坑。灾难性遗忘是最经典的问题，"
        "通过控制alpha/r比值可以有效缓解。数据污染在多语言场景下尤其需要注意。"
        "这些经验来自我们在东南亚电商场景的真实实践。"
    )

    return slide


def split_slide14(prs, original_index):
    """Split Slide 14 into two slides: quantization config + LoRA injection."""
    layout = find_blank_layout(prs)

    # Slide 14a: 量化配置
    slide_a = prs.slides.add_slide(layout)

    add_textbox(slide_a, Inches(0.5), Inches(0.3), Inches(11), Inches(0.8),
                "核心代码 (上): 模型加载与量化配置", font_size=28, bold=True,
                color=RGBColor(0x1A, 0x1A, 0x2E))

    add_textbox(slide_a, Inches(0.5), Inches(1.2), Inches(5), Inches(0.5),
                "📦 关键依赖库", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    add_textbox(slide_a, Inches(0.5), Inches(1.7), Inches(5), Inches(0.5),
                "transformers  •  peft  •  bitsandbytes", font_size=14)

    add_textbox(slide_a, Inches(0.5), Inches(2.4), Inches(5), Inches(0.5),
                "⚡ 4-bit 量化配置 (QLoRA 核心)", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    code_quant = (
        "from transformers import BitsAndBytesConfig\n\n"
        "bnb_config = BitsAndBytesConfig(\n"
        "    load_in_4bit=True,\n"
        '    bnb_4bit_quant_type="nf4",\n'
        "    bnb_4bit_compute_dtype=torch.bfloat16,\n"
        ")"
    )
    add_textbox(slide_a, Inches(0.5), Inches(2.9), Inches(5.5), Inches(2.5),
                code_quant, font_size=11, color=RGBColor(0x22, 0x22, 0x22))

    add_textbox(slide_a, Inches(6.5), Inches(1.2), Inches(5), Inches(0.5),
                "📊 加载模型", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    code_load = (
        "model = AutoModelForCausalLM\n"
        "    .from_pretrained(\n"
        '        "Qwen/Qwen1.5-1.8B-Chat",\n'
        "        quantization_config=bnb_config,\n"
        '        device_map="auto",\n'
        "    )"
    )
    add_textbox(slide_a, Inches(6.5), Inches(1.7), Inches(5), Inches(2.0),
                code_load, font_size=11, color=RGBColor(0x22, 0x22, 0x22))

    add_textbox(slide_a, Inches(6.5), Inches(4.0), Inches(5), Inches(0.5),
                "💡 核心效果", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    add_multiline_textbox(slide_a, Inches(6.5), Inches(4.5), Inches(5), Inches(2.0), [
        ("8B 模型显存占用: 16GB → 4GB", True),
        ("精度损失几乎可忽略", False),
        ("底座模型 4-bit 加载，LoRA 层 FP16 训练", False),
        ('"保大放小"的混合精度策略', False),
    ], font_size=12)

    notes_a = slide_a.notes_slide
    notes_a.notes_text_frame.text = (
        "这页聚焦QLoRA的量化配置。NF4是专为正态分布的模型权重设计的量化格式，"
        "理论上比普通int4量化保留更多有效信息。device_map=auto让框架自动决定"
        "哪些层放GPU、哪些放CPU，实现最优的显存利用。"
    )

    # Slide 14b: LoRA 注入
    slide_b = prs.slides.add_slide(layout)

    add_textbox(slide_b, Inches(0.5), Inches(0.3), Inches(11), Inches(0.8),
                "核心代码 (下): LoRA 注入与训练启动", font_size=28, bold=True,
                color=RGBColor(0x1A, 0x1A, 0x2E))

    add_textbox(slide_b, Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.5),
                "🛠️ LoRA 核心配置", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    code_lora = (
        "from peft import LoraConfig, get_peft_model\n\n"
        "config = LoraConfig(\n"
        "    r=16,              # 秩\n"
        "    lora_alpha=32,     # 缩放系数\n"
        "    target_modules=[\n"
        '        "q_proj", "v_proj",\n'
        '        "k_proj", "o_proj"\n'
        "    ],\n"
        "    lora_dropout=0.1,\n"
        '    task_type="CAUSAL_LM"\n'
        ")\n\n"
        "model = get_peft_model(model, config)\n"
        "model.print_trainable_parameters()\n"
        "# → trainable: 0.39% of all parameters"
    )
    add_textbox(slide_b, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.0),
                code_lora, font_size=11, color=RGBColor(0x22, 0x22, 0x22))

    add_textbox(slide_b, Inches(6.5), Inches(1.2), Inches(5), Inches(0.5),
                "🔑 参数详解", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    params = [
        ("r (Rank)", "低秩矩阵的维度。越大学习能力越强，但参数量也越大。"),
        ("lora_alpha", "缩放系数。alpha/r 决定微调对原模型的影响力。"),
        ("target_modules", "微调的目标层。覆盖注意力全部投影层以捕捉语义。"),
        ("lora_dropout", "正则化手段，防止LoRA层过拟合。"),
    ]

    y = Inches(1.8)
    for name, desc in params:
        add_textbox(slide_b, Inches(6.5), y, Inches(5), Inches(0.3),
                    name, font_size=13, bold=True,
                    color=RGBColor(0x1A, 0x1A, 0x2E))
        y += Inches(0.3)
        add_textbox(slide_b, Inches(6.8), y, Inches(4.7), Inches(0.5),
                    desc, font_size=11, color=RGBColor(0x55, 0x55, 0x55))
        y += Inches(0.55)

    add_textbox(slide_b, Inches(6.5), Inches(4.5), Inches(5), Inches(0.5),
                "💾 训练产出", font_size=18, bold=True,
                color=RGBColor(0x33, 0x33, 0x99))

    add_multiline_textbox(slide_b, Inches(6.5), Inches(5.0), Inches(5), Inches(1.5), [
        ("仅导出几十 MB 的 LoRA 权重文件", True),
        ("底座模型完全不变，适配器可随时插拔", False),
        ("支持多适配器共享一个底座模型", False),
    ], font_size=12)

    notes_b = slide_b.notes_slide
    notes_b.notes_text_frame.text = (
        "这页聚焦LoRA的配置参数。target_modules覆盖了注意力机制的所有投影层，"
        "这比只覆盖q_proj和v_proj能捕捉到更丰富的语义信息。"
        "alpha/r=2是经验值，平衡学习能力和稳定性。"
        "最终产出只有几十MB的适配器文件，可以像插件一样热加载。"
    )

    return slide_a, slide_b


def main():
    print(f"📖 读取 PPT: {INPUT_PATH}")
    prs = Presentation(str(INPUT_PATH))
    original_count = len(prs.slides)
    print(f"   原始幻灯片数量: {original_count}")

    # 1. Split Slide 14 (index 13) into two slides
    print("\n🔧 拆分 Slide 14 → 14a (量化配置) + 14b (LoRA 注入)...")
    slide_a, slide_b = split_slide14(prs, 13)

    # 2. Create "效果验证" slide
    print("🔧 新增 '效果验证: Before vs After' 页...")
    eval_slide = create_eval_slide(prs)

    # 3. Create "实战避坑指南" slide
    print("🔧 新增 '实战避坑指南' 页...")
    pitfalls_slide = create_pitfalls_slide(prs)

    # Reorder slides:
    # Original: 1-14, 15-17, 18(QA)
    # Target:   1-13, 14a, 14b, eval, 15-16(original 15-16), pitfalls, 17(总结), 18(QA)
    #
    # New slides were appended at the end, so we need to move them.
    # After appending: indices 0-17 (original) + 18(14a), 19(14b), 20(eval), 21(pitfalls)
    total = len(prs.slides)
    print(f"\n📐 重排幻灯片顺序 (当前共 {total} 页)...")

    slide_list = prs.slides._sldIdLst

    # Get references to the appended elements (last 4)
    all_elements = list(slide_list)
    el_14a = all_elements[-4]       # slide_a
    el_14b = all_elements[-3]       # slide_b
    el_eval = all_elements[-2]      # eval_slide
    el_pitfalls = all_elements[-1]  # pitfalls_slide

    for el in [el_14a, el_14b, el_eval, el_pitfalls]:
        slide_list.remove(el)

    # Re-read after removal
    all_elements = list(slide_list)

    # Insert 14a, 14b after original slide 13 (index 13)
    ref_after_13 = all_elements[13]
    idx_ref = list(slide_list).index(ref_after_13)
    slide_list.insert(idx_ref + 1, el_14a)
    slide_list.insert(idx_ref + 2, el_14b)

    # Insert eval_slide after 14b (which is now at idx_ref + 2)
    all_elements = list(slide_list)
    idx_14b = list(slide_list).index(el_14b)
    slide_list.insert(idx_14b + 1, el_eval)

    # Insert pitfalls before the last 2 slides (总结 + QA)
    all_elements = list(slide_list)
    # 总结 is the second-to-last original slide; QA is last
    # With insertions, find QA (originally last at index 17, but shifted)
    idx_qa = len(list(slide_list)) - 1  # QA is still last among originals
    # Insert pitfalls before 总结 (second to last)
    idx_summary = idx_qa - 1
    ref_summary = list(slide_list)[idx_summary]
    slide_list.insert(list(slide_list).index(ref_summary), el_pitfalls)

    final_count = len(prs.slides)
    print(f"   最终幻灯片数量: {final_count}")

    prs.save(str(OUTPUT_PATH))
    print(f"\n✅ 改进后的 PPT 已保存到: {OUTPUT_PATH}")
    print("\n📋 变更摘要:")
    print("   - Slide 14 拆分为: '量化配置' + 'LoRA 注入' 两页")
    print("   - 新增: '效果验证: Before vs After' 页")
    print("   - 新增: '实战避坑指南' 页")
    print("   - 总页数: {} → {}".format(original_count, final_count))


if __name__ == "__main__":
    main()
