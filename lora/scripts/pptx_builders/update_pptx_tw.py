"""Update all PPT slides: replace Indonesian references with Taiwanese e-commerce context."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def add_multi(slide, left, top, width, height, lines, size=12, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = color or RGBColor(0x44, 0x44, 0x44)
        p.space_after = Pt(2)
    return tb


def style_cell(cell, val, size=11, bold=False, bg=None, fg=None):
    cell.text = str(val)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        if fg:
            p.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


def clear_slide(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def text_replace_in_slide(slide, replacements):
    """Do text replacement in all shapes of a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            for old, new in replacements.items():
                                if old in run.text:
                                    run.text = run.text.replace(old, new)


def rebuild_slide6_rag(slide):
    """Rebuild Slide 6: LoRA vs RAG — now with Taiwanese slang context."""
    clear_slide(slide)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    RED = RGBColor(0xCC, 0x33, 0x33)
    GREEN = RGBColor(0x22, 0x8B, 0x22)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u7b49\u4e00\u4e0b\uff0c\u70ba\u4ec0\u9ebc\u4e0d\u7528 RAG\uff1f",
             size=28, bold=True, color=DARK)
    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "LoRA \u5fae\u8c03 vs RAG \u68c0\u7d22\u589e\u5f3a\uff1a"
             "\u4e0d\u540c\u95ee\u9898\u7528\u4e0d\u540c\u5de5\u5177",
             size=14, color=GRAY)

    tbl_shape = slide.shapes.add_table(7, 3,
        Inches(0.3), Inches(1.35), Inches(11.4), Inches(2.8))
    tbl = tbl_shape.table
    headers = ["\u5c0d\u6bd4\u7dad\u5ea6",
               "RAG\uff08\u6aa2\u7d22\u589e\u5f37\u751f\u6210\uff09",
               "LoRA\uff08\u4f4e\u79e9\u5fae\u8abf\uff09"]
    rows = [
        ["\u6838\u5fc3\u539f\u7406",
         "\u63a8\u7406\u6642\u6aa2\u7d22\u76f8\u95dc\u6587\u6a94\uff0c\u585e\u9032 Prompt",
         "\u8a13\u7df4\u6642\u628a\u77e5\u8b58\u5beb\u5165\u6a21\u578b\u6b0a\u91cd"],
        ["\u985e\u6bd4",
         "\u201c\u958b\u5377\u8003\u8a66\u201d\u2014\u2014\u5e36\u8457\u8cc7\u6599\u67e5",
         "\u201c\u9589\u5377\u8003\u8a66\u201d\u2014\u2014\u77e5\u8b58\u5df2\u5167\u5316"],
        ["\u64c5\u9577\u5834\u666f",
         "\u4e8b\u5be6\u67e5\u8a62\u3001\u77e5\u8b58\u5eab\u554f\u7b54\u3001\u6587\u4ef6\u641c\u7d22",
         "\u98a8\u683c\u9077\u79fb\u3001\u8a9e\u8a00\u7fd2\u6163\u3001\u9818\u57df\u884c\u70ba\u6a21\u5f0f"],
        ["\u63a8\u7406\u5ef6\u9072",
         "\u8f03\u9ad8\uff08\u9700\u8981\u5148\u6aa2\u7d22\u518d\u751f\u6210\uff09",
         "\u4f4e\uff08\u77e5\u8b58\u5df2\u5728\u6a21\u578b\u5167\uff09"],
        ["\u77e5\u8b58\u66f4\u65b0",
         "\u5373\u6642\u66f4\u65b0\uff08\u6539\u6587\u4ef6\u5373\u53ef\uff09",
         "\u9700\u91cd\u65b0\u8a13\u7df4"],
    ]
    for ci, v in enumerate(headers):
        style_cell(tbl.cell(0, ci), v, size=11, bold=True, bg="333399",
                   fg=RGBColor(0xFF, 0xFF, 0xFF))
    for ri, row in enumerate(rows):
        bg = "F5F5FF" if ri % 2 == 0 else "FFFFFF"
        for ci, v in enumerate(row):
            style_cell(tbl.cell(ri + 1, ci), v, size=10, bg=bg)
    style_cell(tbl.cell(6, 0), rows[4][0], size=10, bg="FFF5E0")
    style_cell(tbl.cell(6, 1), rows[4][1], size=10, bg="E8FFE8", fg=GREEN)
    style_cell(tbl.cell(6, 2), rows[4][2], size=10, bg="FFE8E8", fg=RED)

    add_text(slide, Inches(0.3), Inches(4.3), Inches(11), Inches(0.4),
             "\U0001f3af \u53f0\u7063\u96fb\u5546\u9ed1\u8a71\u5834\u666f\uff1a"
             "\u70ba\u4ec0\u9ebc LoRA \u66f4\u5408\u9069\uff1f",
             size=16, bold=True, color=BLUE)

    add_text(slide, Inches(0.3), Inches(4.75), Inches(5.5), Inches(0.3),
             "\u274c RAG \u7684\u56f0\u5883", size=13, bold=True, color=RED)
    add_multi(slide, Inches(0.3), Inches(5.05), Inches(5.5), Inches(1.8), [
        ("\u300c\u624b\u5200\u4e0b\u55ae\u300d\u300c\u6bcd\u6e6f\u300d\u300cCP\u503c\u300d\u300c\u4f5b\u5fc3\u50f9\u300d"
         "\u9019\u985e\u9ed1\u8a71\u4e0d\u662f\u300c\u4e8b\u5be6\u300d\uff0c\u800c\u662f\u300c\u8a9e\u8a00\u7fd2\u6163\u300d\u3002", False),
        ("", False),
        ("\u2022 \u4f60\u6c92\u6cd5\u6aa2\u7d22\u300c\u5982\u4f55\u5728\u53e5\u5b50\u88e1\u81ea\u7136\u5730\u7528\u6bcd\u6e6f\u300d", False),
        ("\u2022 \u5c31\u7b97\u6aa2\u7d22\u5230\u4e86\uff0c\u6a21\u578b\u4e5f\u4e0d\u6703\u81ea\u7136\u5730\u300c\u8aaa\u300d\u51fa\u4f86", False),
        ("\u2022 \u5c31\u50cf\u7d66\u5916\u570b\u4eba\u4e00\u672c\u53f0\u8a9e\u8fad\u5178\uff0c\u4ed6\u67e5\u5f97\u5230\u4f46\u8aaa\u4e0d\u9053\u5730", False),
    ], size=11, color=GRAY)

    add_text(slide, Inches(6.0), Inches(4.75), Inches(5.5), Inches(0.3),
             "\u2705 LoRA \u7684\u512a\u52e2", size=13, bold=True, color=GREEN)
    add_multi(slide, Inches(6.0), Inches(5.05), Inches(5.5), Inches(1.8), [
        ("\u9ed1\u8a71\u662f\u4e00\u7a2e\u300c\u8a9e\u611f\u300d\uff0c"
         "\u9700\u8981\u5beb\u9032\u6a21\u578b\u7684\u300c\u8a9e\u8a00\u4e2d\u6a1e\u300d\u3002", False),
        ("", False),
        ("\u2022 \u5fae\u8abf\u5f8c\u6a21\u578b\u80fd\u81ea\u7136\u5730\u5728\u4e0a\u4e0b\u6587\u4e2d\u7a7f\u63d2\u9ed1\u8a71", False),
        ("\u2022 \u4e0d\u9700\u8981\u6aa2\u7d22\uff0c\u63a8\u7406\u6642\u96f6\u984d\u5916\u5ef6\u9072", False),
        ("\u2022 \u5c31\u50cf\u8b93\u5916\u570b\u4eba\u5728\u53f0\u7063\u4f4f\u4e09\u500b\u6708\uff0c"
         "\u81ea\u7136\u5c31\u6703\u8aaa\u300c\u8e29\u96f7\u300d\u300c\u79d2\u6bba\u300d", False),
    ], size=11, color=GRAY)

    add_text(slide, Inches(0.3), Inches(6.55), Inches(11.4), Inches(0.5),
             "\U0001f4a1 \u7d50\u8ad6\uff1a"
             "RAG \u89e3\u6c7a\u300c\u6a21\u578b\u4e0d\u77e5\u9053\u300d\uff0c"
             "LoRA \u89e3\u6c7a\u300c\u6a21\u578b\u4e0d\u6703\u8aaa\u300d\u3002"
             "\u5169\u8005\u4e0d\u5c0d\u7acb\uff0c\u8907\u96dc\u5834\u666f\u53ef\u4ee5\u7d44\u5408\u4f7f\u7528"
             "\uff08LoRA \u5b78\u98a8\u683c + RAG \u67e5\u4e8b\u5be6\uff09\u3002",
             size=13, bold=True, color=DARK)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u9019\u9801\u56de\u7b54\u807d\u773e\u6700\u5e38\u554f\u7684\u554f\u984c\uff1a\u70ba\u4ec0\u9ebc\u4e0d\u7528RAG\uff1f"
        "\u53f0\u7063\u96fb\u5546\u9ed1\u8a71\u4e0d\u662f\u4e8b\u5be6\u6027\u77e5\u8b58\uff0c\u800c\u662f\u8a9e\u8a00\u7fd2\u6163\u3002"
    )


def rebuild_slide22_before_after(slide):
    """Rebuild Slide 22: Before vs After with Taiwanese examples."""
    clear_slide(slide)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    RED = RGBColor(0xCC, 0x33, 0x33)
    GREEN = RGBColor(0x22, 0x8B, 0x22)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u6548\u679c\u9a57\u8b49\uff1aBefore vs After",
             size=28, bold=True, color=DARK)
    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "\u540c\u4e00\u500b Prompt\uff0c\u5fae\u8abf\u524d\u5f8c\u7684\u8f38\u51fa\u5c0d\u6bd4",
             size=14, color=GRAY)

    add_text(slide, Inches(0.3), Inches(1.3), Inches(11), Inches(0.4),
             "\U0001f4dd Prompt: \u300c\u4f60\u662f\u53f0\u7063\u96fb\u5546\u7db2\u7d05\uff0c"
             "\u8acb\u7528\u53f0\u7063\u53e3\u8a9e\u98a8\u683c\u63a8\u5ee3\u9019\u6b3e\u884c\u52d5\u96fb\u6e90: "
             "20000mAh\uff0c\u7279\u50f9 NT$499\u300d",
             size=12, bold=True, color=DARK)

    # Before
    add_text(slide, Inches(0.3), Inches(1.9), Inches(5.5), Inches(0.3),
             "\U0001f534 Before\uff08\u901a\u7528\u6a21\u578b\uff09", size=14, bold=True, color=RED)
    add_multi(slide, Inches(0.3), Inches(2.2), Inches(5.5), Inches(1.8), [
        ("\u9019\u6b3e 20000mAh \u884c\u52d5\u96fb\u6e90\u5177\u6709\u5927\u5bb9\u91cf\u3001"
         "\u8f15\u4fbf\u651c\u5e36\u7b49\u7279\u9ede\u3002\u652f\u63f4\u96d9\u5411\u5feb\u5145\u6280\u8853\uff0c"
         "\u53ef\u4ee5\u6eff\u8db3\u60a8\u7684\u65e5\u5e38\u5145\u96fb\u9700\u6c42\u3002"
         "\u7279\u50f9 NT$499\uff0c\u6027\u50f9\u6bd4\u8f03\u9ad8...", False),
        ("", False),
        ("\u26a0\ufe0f \u554f\u984c: \u50cf\u7522\u54c1\u8aaa\u660e\u66f8\uff0c"
         "\u6c92\u6709\u53f0\u7063\u96fb\u5546\u6c1b\u570d", True),
    ], size=11, color=GRAY)

    # After
    add_text(slide, Inches(6.0), Inches(1.9), Inches(5.5), Inches(0.3),
             "\U0001f7e2 After\uff08LoRA \u5fae\u8abf\u5f8c\uff09", size=14, bold=True, color=GREEN)
    add_multi(slide, Inches(6.0), Inches(2.2), Inches(5.5), Inches(1.8), [
        ("\u5bf6\u5011\u9019\u6ce2\u4e0d\u8cb7\u771f\u7684\u6703\u5f8c\u6094\uff01\U0001f525 "
         "\u9019\u9846\u884c\u52d5\u96fb\u6e90 20000mAh \u8d85\u8584\u8d85\u8f15\uff0c"
         "\u5feb\u5145\u8d85\u7d66\u529b\uff01\u624b\u5200\u4e0b\u55ae\u624d $499\uff0c"
         "CP\u503c\u7206\u8868\u6839\u672c\u4f5b\u5fc3\u50f9\uff01"
         "\u514d\u904b\u76f4\u63a5\u7d50\u5e33\uff0c\u79d2\u6bba\u5c31\u6c92\u4e86\uff01\U0001f6d2\u2728", False),
        ("", False),
        ("\u2705 \u81ea\u7136\u7a7f\u63d2\u53f0\u7063\u9ed1\u8a71\uff0c\u50cf\u771f\u4eba\u7db2\u7d05\u5728\u5e36\u8ca8", True),
    ], size=11, color=GRAY)

    # Highlighted slang
    add_text(slide, Inches(0.3), Inches(4.2), Inches(11), Inches(0.4),
             "\U0001f50d \u5fae\u8abf\u5f8c\u81ea\u7136\u51fa\u73fe\u7684\u53f0\u7063\u96fb\u5546\u9ed1\u8a71:",
             size=14, bold=True, color=BLUE)

    tbl = slide.shapes.add_table(2, 6,
        Inches(0.3), Inches(4.6), Inches(11.4), Inches(0.8))
    t = tbl.table
    terms = ["\u624b\u5200\u4e0b\u55ae", "CP\u503c", "\u4f5b\u5fc3\u50f9",
             "\u514d\u904b", "\u79d2\u6bba", "\u7d50\u5e33"]
    meanings = ["rush to buy", "cost-performance", "generous price",
                "free shipping", "instant sellout", "checkout"]
    for ci, v in enumerate(terms):
        style_cell(t.cell(0, ci), v, size=12, bold=True, bg="E8FFE8",
                   fg=GREEN)
    for ci, v in enumerate(meanings):
        style_cell(t.cell(1, ci), v, size=10, bg="F5F5FF")

    add_text(slide, Inches(0.3), Inches(5.7), Inches(11.4), Inches(0.5),
             "\U0001f4a1 \u50c5\u7528 30 \u689d\u6a23\u672c\u8cc7\u6599 + \u5e7e\u5206\u9418\u8a13\u7df4\uff0c"
             "\u6a21\u578b\u5c31\u80fd\u5f9e\u300c\u7522\u54c1\u8aaa\u660e\u66f8\u98a8\u683c\u300d"
             "\u8f49\u8b8a\u70ba\u300c\u53f0\u7063\u7db2\u7d05\u5e36\u8ca8\u98a8\u683c\u300d\uff01",
             size=14, bold=True, color=DARK)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u9019\u662f\u6574\u500b Demo \u7684\u9ad8\u6f6e\u3002"
        "\u5de6\u908a\u662f\u901a\u7528\u6a21\u578b\u7684\u751f\u786c\u8f38\u51fa\uff0c"
        "\u53f3\u908a\u662f\u5fae\u8abf\u5f8c\u81ea\u7136\u5e36\u53f0\u7063\u9ed1\u8a71\u7684\u8f38\u51fa\u3002"
    )


def rebuild_slide24_roadmap(slide):
    """Rebuild Slide 24: data roadmap — focus on TW + cross-strait scenarios."""
    clear_slide(slide)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    ORANGE = RGBColor(0xCC, 0x66, 0x00)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u6578\u64da\u898f\u6a21\u8207\u5834\u666f\u64f4\u5c55\u8def\u7dda\u5716",
             size=28, bold=True, color=DARK)
    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "\u5f9e Demo \u73a9\u5177\u5230\u751f\u7522\u53ef\u7528\uff0c\u6578\u64da\u662f\u6838\u5fc3\u7684\u300c\u6cb9\u9580\u300d",
             size=14, color=GRAY)

    add_text(slide, Inches(0.3), Inches(1.35), Inches(11), Inches(0.35),
             "\U0001f4c8 \u6578\u64da\u91cf\u8def\u7dda\u5716",
             size=16, bold=True, color=BLUE)

    tbl = slide.shapes.add_table(5, 5,
        Inches(0.3), Inches(1.75), Inches(11.4), Inches(2.0))
    t = tbl.table
    headers = ["\u968e\u6bb5", "\u6578\u64da\u91cf",
               "\u9810\u671f\u6548\u679c", "\u8a13\u7df4\u6642\u9593 (1xA100)",
               "\u9069\u7528\u5834\u666f"]
    for ci, v in enumerate(headers):
        style_cell(t.cell(0, ci), v, size=10, bold=True, bg="333399",
                   fg=RGBColor(0xFF, 0xFF, 0xFF))
    rows = [
        ["\U0001f6a9 Demo", "30 \u689d", "\u8ddf\u901a\u6d41\u7a0b\uff0c\u6548\u679c\u96a8\u6a5f", "< 5 \u5206\u9418", "\u6559\u5b78\u6f14\u793a"],
        ["\U0001f680 MVP", "200-500 \u689d", "\u57fa\u672c\u98a8\u683c\u9077\u79fb", "10-30 \u5206\u9418", "\u5167\u90e8\u8a66\u7528"],
        ["\u2705 \u53ef\u7528\u7d1a", "1000-3000 \u689d", "\u7a69\u5b9a\u8f38\u51fa\u9818\u57df\u5167\u5bb9", "1-3 \u5c0f\u6642", "\u5c0f\u7bc4\u570d\u4e0a\u7dda"],
        ["\U0001f31f \u751f\u7522\u7d1a", "5000+ \u689d", "\u98a8\u683c\u81ea\u7136\u3001\u8986\u84cb\u5168\u9762", "3-8 \u5c0f\u6642", "\u5927\u898f\u6a21\u90e8\u7f72"],
    ]
    colors = ["FFF5E0", "E8F0FF", "E8FFE8", "F5E8FF"]
    for ri, (row, bg) in enumerate(zip(rows, colors)):
        for ci, v in enumerate(row):
            style_cell(t.cell(ri + 1, ci), v, size=9, bg=bg)

    add_text(slide, Inches(0.3), Inches(3.9), Inches(11), Inches(0.35),
             "\U0001f30f \u5834\u666f\u64f4\u5c55\u4e09\u7dad\u5ea6",
             size=16, bold=True, color=BLUE)

    # Dimension 1: Markets
    add_text(slide, Inches(0.3), Inches(4.35), Inches(3.6), Inches(0.3),
             "\U0001f30f \u591a\u5e02\u5834", size=13, bold=True, color=ORANGE)
    add_multi(slide, Inches(0.3), Inches(4.65), Inches(3.6), Inches(1.3), [
        ("\u53f0\u7063: \u624b\u5200\u4e0b\u55ae\u3001\u6bcd\u6e6f\u3001CP\u503c (\u7576\u524d)", True),
        ("\u6e2f\u6fb3: \u5289\u4ee5\u9054\u3001\u6397\u5e95\u50f9\u3001\u52c1\u6263", False),
        ("\u5927\u9678: \u7a2e\u8349\u3001\u62d4\u8349\u3001YYDS", False),
        ("\u6771\u5357\u4e9e: cuan\u3001spill\u3001sikat (ID)", False),
        ("\u65e5\u97e9: \u30b3\u30b9\u30d1\u3001\u30dd\u30c1\u308b\u3001\u5272\u5f15 (KR/JP)", False),
    ], size=10, color=GRAY)

    # Dimension 2: Tasks
    add_text(slide, Inches(4.1), Inches(4.35), Inches(3.6), Inches(0.3),
             "\U0001f4dd \u591a\u4efb\u52d9", size=13, bold=True, color=ORANGE)
    add_multi(slide, Inches(4.1), Inches(4.65), Inches(3.6), Inches(1.3), [
        ("\u5546\u54c1\u63a8\u5ee3\u6587\u6848 (\u7576\u524d\u5df2\u5be6\u73fe)", True),
        ("\u5ba2\u670d\u8a71\u8853 (\u7576\u524d\u5df2\u5be6\u73fe)", True),
        ("\u793e\u7fa4\u5a92\u9ad4\u7a2e\u8349\u6587", False),
        ("\u76f4\u64ad\u5e36\u8ca8\u8173\u672c", False),
        ("\u5718\u8cfc\u958b\u5718\u6587\u6848", False),
    ], size=10, color=GRAY)

    # Dimension 3: Formats
    add_text(slide, Inches(7.9), Inches(4.35), Inches(3.6), Inches(0.3),
             "\U0001f4f1 \u591a\u683c\u5f0f", size=13, bold=True, color=ORANGE)
    add_multi(slide, Inches(7.9), Inches(4.65), Inches(3.6), Inches(1.3), [
        ("\u77ed\u6587\u6848: TikTok/Reels (\u226480\u5b57)", False),
        ("\u9577\u6587\u6848: \u90e8\u843d\u683c/\u8a73\u60c5\u9801 (300+\u5b57)", False),
        ("\u5c0d\u8a71\u9ad4: \u5ba2\u670d\u804a\u5929/\u76f4\u64ad\u4e92\u52d5", False),
        ("\u7d50\u69cb\u5316: \u5546\u54c1\u6a19\u984c + \u8ce3\u9ede\u5217\u8868", False),
        ("\u591a\u8f2a\u5c0d\u8a71: \u8aca\u8a62\u8ddf\u9032\u5834\u666f", False),
    ], size=10, color=GRAY)

    add_text(slide, Inches(0.3), Inches(6.2), Inches(11.4), Inches(0.6),
             "\U0001f4a1 \u95dc\u9375\u539f\u5247\uff1a"
             "\u54c1\u8cea > \u6578\u91cf\u3002"
             "100 \u689d\u9ad8\u54c1\u8cea\u6578\u64da > 1000 \u689d\u566a\u97f3\u6578\u64da\u3002"
             "\u6bcf\u500b\u5e02\u5834/\u4efb\u52d9\u5efa\u8b70\u7368\u7acb\u8a13\u7df4\u4e00\u500b LoRA \u9069\u914d\u5668\uff0c"
             "\u63a8\u7406\u6642\u901a\u904e Multi-LoRA \u52d5\u614b\u5207\u63db\u3002",
             size=13, bold=True, color=DARK)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u6f14\u793a\u7528\u768430\u689d\u6578\u64da\u53ea\u662f\u8d77\u9ede\uff0c"
        "\u751f\u7522\u7d1a\u5225\u9700\u89815000+\u689d\u3002"
        "\u5834\u666f\u64f4\u5c55\u4e09\u7dad\u5ea6\u53ef\u4ee5\u6392\u5217\u7d44\u5408\u3002"
    )


def rebuild_slide27_pitfalls(slide):
    """Rebuild Slide 27: pitfalls — replace Indonesian examples with Taiwanese."""
    clear_slide(slide)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    ORANGE = RGBColor(0xCC, 0x66, 0x00)
    GRAY = RGBColor(0x44, 0x44, 0x44)

    add_text(slide, Inches(0.5), Inches(0.25), Inches(11), Inches(0.6),
             "\u5be6\u6230\u907f\u5751\u6307\u5357",
             size=28, bold=True, color=DARK)

    pitfalls = [
        ("\u26a0\ufe0f \u707d\u96e3\u6027\u907a\u5fd8 (Catastrophic Forgetting)",
         "\u6a21\u578b\u5b78\u6703\u65b0\u77e5\u8b58\u4f46\u5fd8\u8a18\u820a\u80fd\u529b\uff0c"
         "\u8868\u73fe\u70ba\u300c\u8907\u8b80\u6a5f\u300d\u6216\u908f\u8f2f\u6df7\u4e82\u3002"
         "\u89e3\u6cd5\uff1a\u964d\u4f4e lora_alpha / r \u6bd4\u503c\uff0c\u63a8\u85a6 alpha = 2r\u3002"),

        ("\u26a0\ufe0f \u6578\u64da\u6c61\u67d3 (Data Contamination)",
         "\u53f0\u7063\u6578\u64da\u6df7\u5165\u5927\u9678\u7528\u8a9e\uff0c"
         "\u6a21\u578b\u7522\u751f\u98a8\u683c\u4ea4\u53c9\u300c\u5e7b\u89ba\u300d"
         "\uff08\u4f8b\u5982\u53f0\u7063\u6587\u6848\u7a81\u7136\u51fa\u73fe\u300c\u89aa\u300d\u300c\u5b9d\u8d1d\u300d\uff09\u3002"
         "\u89e3\u6cd5\uff1a\u56b4\u683c\u6309\u5e02\u5834/\u98a8\u683c\u9694\u96e2\u6578\u64da\u96c6\u3002"),

        ("\u26a0\ufe0f Rank \u9078\u64c7\u4e0d\u7576",
         "r \u904e\u5c0f\u2192\u5b78\u4e0d\u5230\u9818\u57df\u77e5\u8b58\uff1b"
         "r \u904e\u5927\u2192\u904e\u64ec\u5408\u3002"
         "\u7d93\u9a57\u503c\uff1ar=8 \u98a8\u683c\u9077\u79fb\uff0c"
         "r=16~32 \u8907\u96dc\u8a9e\u7fa9\u3002"),

        ("\u26a0\ufe0f \u7f3a\u5c11\u986f\u5f0f\u6a19\u8a3b",
         "\u50c5\u9760\u8a9e\u8a00\u8b93\u6a21\u578b\u81ea\u52d5\u8b58\u5225\u5e02\u5834\uff0c"
         "\u63a8\u7406\u6642\u5bb9\u6613\u300c\u4e32\u53f0\u300d\u3002"
         "\u89e3\u6cd5\uff1a\u5728 instruction \u4e2d\u6a19\u8a3b\u89d2\u8272\u548c\u5e02\u5834"
         "\uff08\u5982\u300c\u4f60\u662f\u53f0\u7063\u96fb\u5546\u7db2\u7d05\u300d\uff09\u3002"),

        ("\u26a0\ufe0f \u5e95\u5ea7\u6a21\u578b\u7d81\u5b9a (Base Model Lock-in)",
         "LoRA \u9069\u914d\u5668\u8207\u5e95\u5ea7\u5f37\u7d81\u5b9a\uff0c"
         "\u63db\u5e95\u5ea7\u6b0a\u91cd\u5c31\u5c0d\u4e0d\u4e0a\u3002"
         "\u89e3\u6cd5\uff1a\u7d71\u4e00\u7ba1\u7406\u5e95\u5ea7\u7248\u672c\uff0c"
         "\u9069\u914d\u5668\u547d\u540d\u5e36\u5e95\u5ea7\u6a19\u8b58"
         "\uff08\u5982 lora-tw-qwen1.5-1.8b\uff09\u3002"),
    ]

    y = Inches(0.95)
    for title, desc in pitfalls:
        tb = slide.shapes.add_textbox(Inches(0.5), y, Inches(11), Inches(0.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        y += Inches(0.35)

        tb2 = slide.shapes.add_textbox(Inches(0.8), y, Inches(10.5), Inches(0.55))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY
        y += Inches(0.7)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u6578\u64da\u6c61\u67d3\u5728\u53f0\u7063\u5834\u666f\u4e0b\u7279\u5225\u5e38\u898b\uff0c"
        "\u56e0\u70ba\u53f0\u7063\u548c\u5927\u9678\u90fd\u662f\u4e2d\u6587\uff0c"
        "\u6a21\u578b\u5f88\u5bb9\u6613\u6df7\u6dc6\u5169\u908a\u7684\u7528\u8a9e\u98a8\u683c\u3002"
    )


def main():
    prs = Presentation("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")

    # Simple text replacements for slides that just need keyword swaps
    common_replacements = {
        "\u4e1c\u5357\u4e9a\u7535\u5546\u63a8\u5ba2": "\u53f0\u7063\u96fb\u5546\u7db2\u7d05",
        "\u4e1c\u5357\u4e9a\u7535\u5546\u4e13\u5bb6": "\u53f0\u7063\u96fb\u5546\u5c08\u5bb6",
        "\u6253\u9020\u4e1c\u5357\u4e9a": "\u6253\u9020\u53f0\u7063",
        "\u4e1c\u5357\u4e9a": "\u53f0\u7063",
        "\u5370\u5c3c\u672c\u5730": "\u53f0\u7063\u672c\u5730",
        "\u5370\u5c3c\u8d44\u6df1\u7535\u5546\u63a8\u5ba2": "\u53f0\u7063\u8cc7\u6df1\u96fb\u5546\u7db2\u7d05",
        "\u5370\u5c3c\u7535\u5546": "\u53f0\u7063\u96fb\u5546",
        "\u5370\u5c3c": "\u53f0\u7063",
        "Bestie/Spill/Checkout/Gratis Ongkir": "\u624b\u5200\u4e0b\u55ae/CP\u503c/\u514d\u904b/\u79d2\u6bba",
        "100k IDR": "NT$499",
        "Powerbank 20000mAh, \u4ec5\u552e 100k IDR": "\u884c\u52d5\u96fb\u6e90 20000mAh, \u7279\u50f9 NT$499",
    }

    # Apply text replacements to all slides first
    for slide in prs.slides:
        text_replace_in_slide(slide, common_replacements)

    # Now rebuild specific slides that need complete overhaul
    slides = list(prs.slides)

    # Find slides by title keywords
    slide_map = {}
    for i, slide in enumerate(slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if "\u7b49\u4e00\u4e0b" in t and "RAG" in t:
                    slide_map["rag"] = i
                elif "\u6548\u679c\u9a8c\u8bc1" in t or "\u6548\u679c\u9a57\u8b49" in t:
                    slide_map["before_after"] = i
                elif "\u6570\u636e\u89c4\u6a21" in t or "\u6578\u64da\u898f\u6a21" in t:
                    slide_map["roadmap"] = i
                elif "\u907f\u5751\u6307\u5357" in t or "\u907f\u5751\u6307\u5357" in t:
                    slide_map["pitfalls"] = i
                break

    print("Slides to rebuild:", slide_map)

    if "rag" in slide_map:
        print(f"  Rebuilding Slide {slide_map['rag']+1}: RAG vs LoRA")
        rebuild_slide6_rag(slides[slide_map["rag"]])

    if "before_after" in slide_map:
        print(f"  Rebuilding Slide {slide_map['before_after']+1}: Before vs After")
        rebuild_slide22_before_after(slides[slide_map["before_after"]])

    if "roadmap" in slide_map:
        print(f"  Rebuilding Slide {slide_map['roadmap']+1}: Data Roadmap")
        rebuild_slide24_roadmap(slides[slide_map["roadmap"]])

    if "pitfalls" in slide_map:
        print(f"  Rebuilding Slide {slide_map['pitfalls']+1}: Pitfalls")
        rebuild_slide27_pitfalls(slides[slide_map["pitfalls"]])

    # Print final order
    print(f"\nTotal: {len(prs.slides)} slides\n")
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    print(f"  Slide {i+1:2d}: {t}")
                    break

    prs.save("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    print("\nDone!")


if __name__ == "__main__":
    main()
