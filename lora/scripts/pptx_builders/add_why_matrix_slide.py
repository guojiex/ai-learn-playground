"""Add a slide explaining why LLMs are fundamentally matrix multiplications, before Slide 9."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def style_cell(cell, val, size=12, bold=False, bg=None, fg=None):
    cell.text = str(val)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        if fg:
            p.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


def main():
    prs = Presentation("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)

    # ===== Title =====
    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u524d\u7f6e\u77e5\u8bc6\uff1a\u4e3a\u4ec0\u4e48\u5927\u6a21\u578b = \u4e00\u5806\u77e9\u9635\uff1f",
             size=28, bold=True, color=DARK)

    add_text(slide, Inches(0.5), Inches(0.85), Inches(11), Inches(0.4),
             "\u7406\u89e3\u8fd9\u4e00\u70b9\uff0c\u540e\u9762\u7684 LoRA \u516c\u5f0f\u5c31\u4e0d\u518d\u62bd\u8c61",
             size=14, color=GRAY)

    # ===== Section 1: Single neuron layer =====
    add_text(slide, Inches(0.3), Inches(1.4), Inches(11), Inches(0.4),
             "\U0001f9e0 \u4e00\u4e2a\u795e\u7ecf\u7f51\u7edc\u5c42\u7684\u672c\u8d28",
             size=16, bold=True, color=BLUE)

    add_text(slide, Inches(0.5), Inches(1.85), Inches(5.5), Inches(0.5),
             "\u8f93\u51fa = \u6fc0\u6d3b\u51fd\u6570( W \u00b7 \u8f93\u5165 + \u504f\u7f6e )",
             size=18, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(6.2), Inches(1.8), Inches(5.5), Inches(0.7),
             "\u6743\u91cd\u77e9\u9635 W \u5c31\u662f\u8fd9\u4e00\u5c42\u201c\u5b66\u5230\u7684\u77e5\u8bc6\u201d\u3002\n"
             "\u4e00\u4e2a d_out \u00d7 d_in \u7684\u77e9\u9635\uff0c\u628a\u8f93\u5165\u5411\u91cf\u6620\u5c04\u5230\u8f93\u51fa\u7a7a\u95f4\u3002",
             size=12, color=GRAY)

    # ===== Section 2: Transformer table =====
    add_text(slide, Inches(0.3), Inches(2.65), Inches(11), Inches(0.4),
             "\U0001f916 Transformer \u91cc\u5230\u5904\u90fd\u662f\u77e9\u9635\u4e58\u6cd5",
             size=16, bold=True, color=BLUE)

    tbl_shape = slide.shapes.add_table(6, 3,
        Inches(0.5), Inches(3.1), Inches(7.5), Inches(2.2))
    tbl = tbl_shape.table

    headers = ["\u7ec4\u4ef6", "\u5bf9\u5e94\u7684\u77e9\u9635", "\u4f5c\u7528"]
    rows = [
        ["Attention Q/K/V \u6295\u5f71", "W_Q, W_K, W_V", "\u628a\u8f93\u5165\u53d8\u6210 Query/Key/Value"],
        ["\u6ce8\u610f\u529b\u8f93\u51fa\u6295\u5f71", "W_O", "\u628a\u6ce8\u610f\u529b\u7ed3\u679c\u6620\u5c04\u56de\u539f\u7a7a\u95f4"],
        ["FFN / MLP \u5c42", "W_up, W_down", "\u4e24\u5c42\u7ebf\u6027\u53d8\u6362 + \u6fc0\u6d3b\u51fd\u6570"],
        ["Embedding \u5c42", "W_embed", "\u628a token ID \u6620\u5c04\u6210\u5411\u91cf"],
        ["LM Head \u8f93\u51fa\u5c42", "W_head", "\u628a\u5411\u91cf\u6620\u5c04\u56de\u8bcd\u8868\u6982\u7387"],
    ]

    for ci, v in enumerate(headers):
        style_cell(tbl.cell(0, ci), v, size=11, bold=True,
                   bg="333399", fg=RGBColor(0xFF, 0xFF, 0xFF))
    for ri, row in enumerate(rows):
        bg_color = "F5F5FF" if ri % 2 == 0 else "FFFFFF"
        for ci, v in enumerate(row):
            style_cell(tbl.cell(ri + 1, ci), v, size=11, bg=bg_color)

    # Highlight note
    add_text(slide, Inches(8.3), Inches(3.1), Inches(3.2), Inches(2.2),
             "\U0001f4a1 \u4e00\u4e2a 1.8B \u53c2\u6570\u7684\u6a21\u578b\uff0c"
             "\u672c\u8d28\u4e0a\u5c31\u662f\u51e0\u767e\u4e2a\u8fd9\u6837\u7684\u77e9\u9635\u5806\u53e0\u5728\u4e00\u8d77\u3002\n\n"
             "\u6240\u8c13\u201c18 \u4ebf\u53c2\u6570\u201d\uff0c\u5c31\u662f\u6240\u6709\u77e9\u9635\u91cc\u6570\u5b57\u7684\u603b\u6570\u3002",
             size=13, bold=True, color=DARK)

    # ===== Section 3: Connection to LoRA =====
    add_text(slide, Inches(0.3), Inches(5.5), Inches(11), Inches(0.4),
             "\U0001f517 \u6240\u4ee5 LoRA \u7684\u516c\u5f0f\u624d\u6210\u7acb",
             size=16, bold=True, color=BLUE)

    add_text(slide, Inches(0.5), Inches(5.9), Inches(11), Inches(0.5),
             "h = W\u2080 \u00b7 x  +  B \u00b7 A \u00b7 x",
             size=20, bold=True, color=DARK, align=PP_ALIGN.LEFT)

    add_text(slide, Inches(0.5), Inches(6.35), Inches(11), Inches(0.7),
             "\u8fd9\u4e0d\u662f\u6bd4\u55bb\uff0c\u662f\u5b57\u9762\u610f\u4e49\uff1a"
             "\u5728\u67d0\u4e00\u5c42\u7684\u77e9\u9635\u4e58\u6cd5\u65c1\u8fb9\uff0c"
             "\u52a0\u4e00\u6761\u201c\u65c1\u8def\u201d(B\u00b7A)\uff0c"
             "\u5bf9\u8be5\u5c42\u8f93\u51fa\u505a\u5c0f\u4fee\u6b63\u3002"
             "\u8bba\u6587\u901a\u5e38\u9009\u62e9\u5bf9 Q \u548c V \u6295\u5f71\u77e9\u9635\u52a0\u65c1\u8def\uff0c"
             "\u56e0\u4e3a\u5b83\u4eec\u5bf9\u4efb\u52a1\u9002\u914d\u6700\u654f\u611f\u3002",
             size=12, color=GRAY)

    # Speaker notes
    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u8fd9\u9875\u662f\u4e3a\u540e\u9762\u7684\u6570\u5b66\u516c\u5f0f\u505a\u94fa\u57ab\u3002"
        "\u5f88\u591a\u542c\u4f17\u770b\u5230 h=Wx \u4f1a\u89c9\u5f97\u62bd\u8c61\uff0c"
        "\u5148\u8ba9\u4ed6\u4eec\u7406\u89e3\u201c\u5927\u6a21\u578b\u5c31\u662f\u4e00\u5806\u77e9\u9635\u201d\u8fd9\u4e2a\u524d\u63d0\uff0c"
        "\u540e\u9762\u7684\u79e9\u5206\u89e3\u3001\u53c2\u6570\u538b\u7f29\u7b49\u5c31\u6c34\u5230\u6e20\u6210\u4e86\u3002"
        "\u53ef\u4ee5\u7c7b\u6bd4\uff1a\u6a21\u578b\u50cf\u4e00\u4e2a\u5de8\u5927\u7684Excel\u8868\u683c\uff0c"
        "\u6bcf\u4e00\u5c42\u5c31\u662f\u4e00\u4e2a\u5de5\u4f5c\u8868\uff0c\u91cc\u9762\u586b\u6ee1\u4e86\u6570\u5b57\u3002"
    )

    # Move to before Slide 9 (index 8), i.e. after Slide 8 (index 7)
    slide_list = prs.slides._sldIdLst
    all_els = list(slide_list)
    new_el = all_els[-1]
    slide_list.remove(new_el)

    all_els = list(slide_list)
    ref = all_els[7]  # Slide 8: "破局之道"
    idx = list(slide_list).index(ref)
    slide_list.insert(idx + 1, new_el)

    print(f"Total: {len(prs.slides)} slides\n")
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    print(f"  Slide {i+1:2d}: {t}")
                    break

    prs.save("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    print("\nDone!")


if __name__ == "__main__":
    main()
