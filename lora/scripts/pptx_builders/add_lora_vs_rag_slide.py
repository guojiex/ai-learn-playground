"""Add a LoRA vs RAG comparison slide after Slide 5 (本次分享目标)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def add_multi(slide, left, top, width, height, lines, size=12, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = clr if clr else (color or RGBColor(0x44, 0x44, 0x44))
        p.space_after = Pt(2)
    return tb


def style_cell(cell, val, size=11, bold=False, bg=None, fg=None):
    cell.text = str(val)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        if fg:
            p.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


def main():
    prs = Presentation("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    GREEN = RGBColor(0x22, 0x8B, 0x22)
    RED = RGBColor(0xCC, 0x33, 0x33)
    ORANGE = RGBColor(0xCC, 0x66, 0x00)

    # ===== Title =====
    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u7b49\u4e00\u4e0b\uff0c\u4e3a\u4ec0\u4e48\u4e0d\u7528 RAG\uff1f",
             size=28, bold=True, color=DARK)

    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "LoRA \u5fae\u8c03 vs RAG \u68c0\u7d22\u589e\u5f3a\uff1a"
             "\u4e0d\u540c\u95ee\u9898\u7528\u4e0d\u540c\u5de5\u5177",
             size=14, color=GRAY)

    # ===== Comparison table =====
    tbl_shape = slide.shapes.add_table(7, 3,
        Inches(0.3), Inches(1.35), Inches(11.4), Inches(2.8))
    tbl = tbl_shape.table

    headers = ["\u5bf9\u6bd4\u7ef4\u5ea6",
               "RAG\uff08\u68c0\u7d22\u589e\u5f3a\u751f\u6210\uff09",
               "LoRA\uff08\u4f4e\u79e9\u5fae\u8c03\uff09"]
    rows = [
        ["\u6838\u5fc3\u539f\u7406",
         "\u63a8\u7406\u65f6\u68c0\u7d22\u76f8\u5173\u6587\u6863\uff0c\u585e\u8fdb Prompt",
         "\u8bad\u7ec3\u65f6\u628a\u77e5\u8bc6\u5199\u5165\u6a21\u578b\u6743\u91cd"],
        ["\u7c7b\u6bd4",
         "\u201c\u5f00\u5377\u8003\u8bd5\u201d\u2014\u2014\u5e26\u7740\u8d44\u6599\u67e5",
         "\u201c\u95ed\u5377\u8003\u8bd5\u201d\u2014\u2014\u77e5\u8bc6\u5df2\u5185\u5316"],
        ["\u64c5\u957f\u573a\u666f",
         "\u4e8b\u5b9e\u67e5\u8be2\u3001\u77e5\u8bc6\u5e93\u95ee\u7b54\u3001\u6587\u6863\u641c\u7d22",
         "\u98ce\u683c\u8fc1\u79fb\u3001\u8bed\u8a00\u4e60\u60ef\u3001\u9886\u57df\u884c\u4e3a\u6a21\u5f0f"],
        ["\u63a8\u7406\u5ef6\u8fdf",
         "\u8f83\u9ad8\uff08\u9700\u8981\u5148\u68c0\u7d22\u518d\u751f\u6210\uff09",
         "\u4f4e\uff08\u77e5\u8bc6\u5df2\u5728\u6a21\u578b\u5185\uff09"],
        ["\u77e5\u8bc6\u66f4\u65b0",
         "\u5b9e\u65f6\u66f4\u65b0\uff08\u6539\u6587\u6863\u5373\u53ef\uff09",
         "\u9700\u91cd\u65b0\u8bad\u7ec3"],
    ]

    for ci, v in enumerate(headers):
        style_cell(tbl.cell(0, ci), v, size=11, bold=True,
                   bg="333399", fg=RGBColor(0xFF, 0xFF, 0xFF))
    for ri, row in enumerate(rows):
        bg = "F5F5FF" if ri % 2 == 0 else "FFFFFF"
        for ci, v in enumerate(row):
            style_cell(tbl.cell(ri + 1, ci), v, size=10, bg=bg)
    # Last row highlight
    style_cell(tbl.cell(6, 0), rows[4][0], size=10, bg="FFF5E0")
    style_cell(tbl.cell(6, 1), rows[4][1], size=10, bg="E8FFE8", fg=GREEN)
    style_cell(tbl.cell(6, 2), rows[4][2], size=10, bg="FFE8E8", fg=RED)

    # ===== Key example =====
    add_text(slide, Inches(0.3), Inches(4.3), Inches(11), Inches(0.4),
             "\U0001f3af \u4e1c\u5357\u4e9a\u9ed1\u8bdd\u573a\u666f\uff1a\u4e3a\u4ec0\u4e48 LoRA \u66f4\u5408\u9002\uff1f",
             size=16, bold=True, color=BLUE)

    # RAG side
    add_text(slide, Inches(0.3), Inches(4.75), Inches(5.5), Inches(0.3),
             "\u274c RAG \u7684\u56f0\u5883", size=13, bold=True, color=RED)
    add_multi(slide, Inches(0.3), Inches(5.05), Inches(5.5), Inches(1.8), [
        ("\u201ccuan\u201d\u201cFYP\u201d\u201cspill\u201d\u201creceh\u201d "
         "\u8fd9\u7c7b\u9ed1\u8bdd\u4e0d\u662f\u201c\u4e8b\u5b9e\u201d\uff0c"
         "\u800c\u662f\u201c\u8bed\u8a00\u4e60\u60ef\u201d\u3002", False, GRAY),
        ("", False, GRAY),
        ("\u2022 \u4f60\u6ca1\u6cd5\u68c0\u7d22\u201c\u5982\u4f55\u5728\u53e5\u5b50\u91cc\u81ea\u7136\u5730\u7528 cuan\u201d",
         False, GRAY),
        ("\u2022 \u5c31\u7b97\u68c0\u7d22\u5230\u4e86\uff0c\u6a21\u578b\u4e5f\u4e0d\u4f1a\u81ea\u7136\u5730"
         "\u201c\u8bf4\u201d\u51fa\u6765", False, GRAY),
        ("\u2022 \u8fd9\u5c31\u50cf\u7ed9\u8001\u5916\u4e00\u672c\u4e2d\u6587\u8bcd\u5178\uff0c"
         "\u4ed6\u67e5\u5f97\u5230\u4f46\u8bf4\u4e0d\u5730\u9053", False, GRAY),
    ], size=11, color=GRAY)

    # LoRA side
    add_text(slide, Inches(6.0), Inches(4.75), Inches(5.5), Inches(0.3),
             "\u2705 LoRA \u7684\u4f18\u52bf", size=13, bold=True, color=GREEN)
    add_multi(slide, Inches(6.0), Inches(5.05), Inches(5.5), Inches(1.8), [
        ("\u9ed1\u8bdd\u662f\u4e00\u79cd\u201c\u8bed\u611f\u201d\uff0c"
         "\u9700\u8981\u5199\u8fdb\u6a21\u578b\u7684\u201c\u8bed\u8a00\u4e2d\u67a2\u201d\u3002", False, GRAY),
        ("", False, GRAY),
        ("\u2022 \u5fae\u8c03\u540e\u6a21\u578b\u80fd\u81ea\u7136\u5730\u5728\u4e0a\u4e0b\u6587\u4e2d"
         "\u7a7f\u63d2\u9ed1\u8bdd", False, GRAY),
        ("\u2022 \u4e0d\u9700\u8981\u68c0\u7d22\uff0c\u63a8\u7406\u65f6\u96f6\u989d\u5916\u5ef6\u8fdf",
         False, GRAY),
        ("\u2022 \u8fd9\u5c31\u50cf\u8ba9\u8001\u5916\u5728\u4e2d\u56fd\u4f4f\u4e09\u4e2a\u6708\uff0c"
         "\u81ea\u7136\u5c31\u4f1a\u8bf4\u201c\u7275\u201d\u201c\u6253call\u201d", False, GRAY),
    ], size=11, color=GRAY)

    # Bottom punchline
    add_text(slide, Inches(0.3), Inches(6.55), Inches(11.4), Inches(0.5),
             "\U0001f4a1 \u7ed3\u8bba\uff1a"
             "RAG \u89e3\u51b3\u201c\u6a21\u578b\u4e0d\u77e5\u9053\u201d\uff0c"
             "LoRA \u89e3\u51b3\u201c\u6a21\u578b\u4e0d\u4f1a\u8bf4\u201d\u3002"
             "\u4e24\u8005\u4e0d\u5bf9\u7acb\uff0c\u590d\u6742\u573a\u666f\u53ef\u4ee5\u7ec4\u5408\u4f7f\u7528"
             "\uff08LoRA \u5b66\u98ce\u683c + RAG \u67e5\u4e8b\u5b9e\uff09\u3002",
             size=13, bold=True, color=DARK)

    # Speaker notes
    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u8fd9\u9875\u56de\u7b54\u542c\u4f17\u6700\u5e38\u95ee\u7684\u95ee\u9898\uff1a"
        "\u4e3a\u4ec0\u4e48\u4e0d\u7528RAG\uff1f"
        "\u6838\u5fc3\u533a\u522b\uff1aRAG\u662f\u201c\u5f00\u5377\u8003\u8bd5\u201d\uff0c"
        "LoRA\u662f\u201c\u95ed\u5377\u8003\u8bd5\u201d\u3002"
        "\u4e1c\u5357\u4e9a\u9ed1\u8bdd\u4e0d\u662f\u4e8b\u5b9e\u6027\u77e5\u8bc6\uff0c"
        "\u800c\u662f\u4e00\u79cd\u8bed\u8a00\u4e60\u60ef\u548c\u98ce\u683c\uff0c"
        "\u4f60\u6ca1\u6cd5\u201c\u68c0\u7d22\u201d\u4e00\u4e2a\u8bed\u611f\u3002"
        "\u4f46\u5982\u679c\u9700\u8981\u67e5\u8be2\u5546\u54c1\u4ef7\u683c\u3001\u4fc3\u9500\u653f\u7b56\u7b49"
        "\u5b9e\u65f6\u53d8\u52a8\u7684\u4fe1\u606f\uff0cRAG\u66f4\u5408\u9002\u3002"
        "\u6700\u4f73\u5b9e\u8df5\u662f\u7ec4\u5408\u4f7f\u7528\uff1a"
        "LoRA\u5b66\u8bed\u8a00\u98ce\u683c + RAG\u67e5\u5b9e\u65f6\u6570\u636e\u3002"
    )

    # Move to after Slide 5 (index 4, "本次分享目标")
    slide_list = prs.slides._sldIdLst
    all_els = list(slide_list)
    new_el = all_els[-1]
    slide_list.remove(new_el)

    all_els = list(slide_list)
    ref = all_els[4]  # Slide 5
    idx = list(slide_list).index(ref)
    slide_list.insert(idx + 1, new_el)

    print(f"Total: {len(prs.slides)} slides\n")
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    print(f"  Slide {i+1:2d}: {t}")
                    break

    prs.save("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    print("\nDone!")


if __name__ == "__main__":
    main()
