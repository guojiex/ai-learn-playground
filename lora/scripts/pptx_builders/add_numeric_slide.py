"""Add a slide with a concrete numerical rank-decomposition example after Slide 9."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def style_cell(cell, val, size=14, bold=False, bg=None, fg=None):
    cell.text = str(val)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        if fg:
            p.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


def add_matrix_table(slide, left, top, data, cell_w, cell_h,
                     size=14, bold=True, bg=None, fg=None):
    rows = len(data)
    cols = len(data[0])
    w = Inches(cell_w * cols)
    h = Inches(cell_h * rows)
    shape = slide.shapes.add_table(rows, cols, left, top, w, h)
    tbl = shape.table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            style_cell(tbl.cell(r, c), val, size=size, bold=bold, bg=bg, fg=fg)
    return shape


def main():
    prs = Presentation("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    layout = find_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    RED = RGBColor(0xCC, 0x33, 0x33)
    GREEN = RGBColor(0x22, 0x8B, 0x22)

    # ===== Title =====
    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u6570\u5b57\u5b9e\u4f8b\uff1a\u4eb2\u624b\u7b97\u4e00\u6b21\u79e9\u5206\u89e3",
             size=28, bold=True, color=DARK)

    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "\u7528\u6700\u7b80\u5355\u7684\u6570\u5b57\u7406\u89e3 \u0394W \u2248 B \u00d7 A",
             size=14, color=GRAY)

    # ===== Section 1: Rank-1 example =====
    add_text(slide, Inches(0.3), Inches(1.3), Inches(11), Inches(0.4),
             "\U0001f9ee Rank-1 \u5206\u89e3\u793a\u4f8b\uff1a"
             "\u4e00\u4e2a 3\u00d73 \u77e9\u9635\u53ef\u4ee5\u88ab\u62c6\u6210\u4e24\u4e2a\u5c0f\u5411\u91cf\u7684\u4e58\u79ef",
             size=15, bold=True, color=BLUE)

    # Matrix B (3x1) - column vector
    mat_y = Inches(1.85)

    add_text(slide, Inches(0.5), Inches(1.75), Inches(1.4), Inches(0.3),
             "B (3\u00d71)", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_matrix_table(slide,
                     Inches(0.7), mat_y + Inches(0.25),
                     [[2], [1], [3]],
                     cell_w=0.55, cell_h=0.32,
                     size=16, bold=True, bg="FFE0E0", fg=RED)

    # Multiplication sign
    add_text(slide, Inches(1.5), mat_y + Inches(0.3), Inches(0.4), Inches(0.4),
             "\u00d7", size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Matrix A (1x3) - row vector
    add_text(slide, Inches(1.9), Inches(1.75), Inches(2.2), Inches(0.3),
             "A (1\u00d73)", size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_matrix_table(slide,
                     Inches(1.9), mat_y + Inches(0.25),
                     [[1, 2, 1]],
                     cell_w=0.55, cell_h=0.32,
                     size=16, bold=True, bg="E0FFE0", fg=GREEN)

    # Equals sign
    add_text(slide, Inches(3.6), mat_y + Inches(0.3), Inches(0.4), Inches(0.4),
             "=", size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Result matrix ΔW (3x3)
    add_text(slide, Inches(4.0), Inches(1.75), Inches(2.2), Inches(0.3),
             "\u0394W (3\u00d73)", size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_matrix_table(slide,
                     Inches(4.0), mat_y + Inches(0.25),
                     [[2, 4, 2],
                      [1, 2, 1],
                      [3, 6, 3]],
                     cell_w=0.55, cell_h=0.32,
                     size=16, bold=True, bg="E0E0FF", fg=BLUE)

    # Calculation hints
    add_text(slide, Inches(6.2), mat_y + Inches(0.05), Inches(5.0), Inches(1.2),
             "\u2190 \u7b2c1\u884c = 2 \u00d7 [1, 2, 1] = [2, 4, 2]\n"
             "\u2190 \u7b2c2\u884c = 1 \u00d7 [1, 2, 1] = [1, 2, 1]\n"
             "\u2190 \u7b2c3\u884c = 3 \u00d7 [1, 2, 1] = [3, 6, 3]",
             size=12, color=GRAY)

    # ===== Section 2: Parameter count =====
    sec2_y = Inches(3.4)
    add_text(slide, Inches(0.3), sec2_y, Inches(11), Inches(0.4),
             "\U0001f4ca \u53c2\u6570\u8ba1\u6570\u5bf9\u6bd4",
             size=15, bold=True, color=BLUE)

    # Small table for comparison
    tbl_shape = slide.shapes.add_table(3, 4,
        Inches(0.5), sec2_y + Inches(0.4), Inches(7.5), Inches(1.1))
    tbl = tbl_shape.table

    headers = ["", "\u77e9\u9635\u5c3a\u5bf8",
               "\u53c2\u6570\u91cf", "\u538b\u7f29\u6bd4"]
    row1 = ["\u5168\u91cf\u5b58\u50a8 \u0394W",
            "3 \u00d7 3", "9", "\u2014"]
    row2 = ["\u79e9\u5206\u89e3 B\u00d7A",
            "3\u00d71 + 1\u00d73", "3 + 3 = 6", "33% \u2193"]

    for ci, v in enumerate(headers):
        style_cell(tbl.cell(0, ci), v, size=11, bold=True, bg="333399", fg=RGBColor(0xFF, 0xFF, 0xFF))
    for ci, v in enumerate(row1):
        style_cell(tbl.cell(1, ci), v, size=11, bg="F5F5FF")
    for ci, v in enumerate(row2):
        style_cell(tbl.cell(2, ci), v, size=11, bg="E8FFE8")

    # ===== Section 3: Scale up! =====
    sec3_y = Inches(5.0)
    add_text(slide, Inches(0.3), sec3_y, Inches(11), Inches(0.4),
             "\U0001f680 \u653e\u5927\u5230\u771f\u5b9e\u6a21\u578b\u7ef4\u5ea6"
             "\uff08d=2048\uff0cr=16\uff09",
             size=15, bold=True, color=BLUE)

    tbl2_shape = slide.shapes.add_table(3, 4,
        Inches(0.5), sec3_y + Inches(0.4), Inches(7.5), Inches(1.1))
    tbl2 = tbl2_shape.table

    h2 = ["", "\u77e9\u9635\u5c3a\u5bf8",
          "\u53c2\u6570\u91cf", "\u538b\u7f29\u6bd4"]
    r1 = ["\u5168\u91cf\u5b58\u50a8 \u0394W",
          "2048 \u00d7 2048", "4,194,304", "\u2014"]
    r2 = ["\u79e9\u5206\u89e3 B\u00d7A (r=16)",
          "2048\u00d716 + 16\u00d72048",
          "32,768 + 32,768 = 65,536",
          "98.4% \u2193 !"]

    for ci, v in enumerate(h2):
        style_cell(tbl2.cell(0, ci), v, size=11, bold=True, bg="333399", fg=RGBColor(0xFF, 0xFF, 0xFF))
    for ci, v in enumerate(r1):
        style_cell(tbl2.cell(1, ci), v, size=11, bg="FFF5F5")
    for ci, v in enumerate(r2):
        style_cell(tbl2.cell(2, ci), v, size=11, bold=True, bg="E8FFE8")

    # Punchline
    add_text(slide, Inches(0.5), Inches(6.55), Inches(11), Inches(0.4),
             "\U0001f4a1 \u540c\u6837\u7684\u6570\u5b66\u539f\u7406\uff0c"
             "\u5c0f\u77e9\u9635\u7701 33%\uff0c"
             "\u5927\u77e9\u9635\u7701 98.4%"
             " \u2014\u2014 \u7ef4\u5ea6\u8d8a\u9ad8\uff0c\u4f4e\u79e9\u5206\u89e3\u8d8a\u5212\u7b97\uff01",
             size=14, bold=True, color=DARK)

    # Speaker notes
    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u8fd9\u9875\u7528\u771f\u5b9e\u6570\u5b57\u6f14\u793a\u79e9\u5206\u89e3\u7684\u8ba1\u7b97\u8fc7\u7a0b\u3002"
        "\u5148\u7528\u4e00\u4e2a\u6781\u7b80\u76843\u00d73\u4f8b\u5b50\u8ba9\u542c\u4f17\u7406\u89e3B\u00d7A\u600e\u4e48\u5f97\u5230\u0394W\uff0c"
        "\u7136\u540e\u653e\u5927\u52302048\u7ef4\u5ea6\u5c55\u793a\u538b\u7f29\u6bd4\u4ece33%\u8df3\u523098.4%\u3002"
        "\u8fd9\u89e3\u91ca\u4e86\u4e3a\u4ec0\u4e48\u201c\u4f4e\u79e9\u201d\u5728\u5927\u6a21\u578b\u573a\u666f\u4e0b\u7279\u522b\u6709\u6548\u3002"
    )

    # Move this slide to after Slide 9 (index 8)
    slide_list = prs.slides._sldIdLst
    all_els = list(slide_list)
    new_el = all_els[-1]
    slide_list.remove(new_el)

    all_els = list(slide_list)
    ref = all_els[8]  # Slide 9
    idx = list(slide_list).index(ref)
    slide_list.insert(idx + 1, new_el)

    # Print final order
    print(f"Total: {len(prs.slides)} slides\n")
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    print(f"  Slide {i+1:2d}: {t}")
                    break

    prs.save("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    print("\nDone!")


if __name__ == "__main__":
    main()
