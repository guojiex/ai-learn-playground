"""Add two slides: evaluation framework + data scaling roadmap, after Slide 22 (效果验证)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def add_multi(slide, left, top, width, height, lines, size=12, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = color or RGBColor(0x44, 0x44, 0x44)
        p.space_after = Pt(2)
    return tb


def style_cell(cell, val, size=11, bold=False, bg=None, fg=None):
    cell.text = str(val)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        if fg:
            p.font.color.rgb = fg
    if bg:
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": bg})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


def create_eval_framework_slide(prs, layout):
    """Slide: evaluation framework for fine-tuned models."""
    slide = prs.slides.add_slide(layout)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    GREEN = RGBColor(0x22, 0x8B, 0x22)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u5982\u4f55\u8bc4\u6d4b\u5fae\u8c03\u6548\u679c\uff1f",
             size=28, bold=True, color=DARK)
    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "\u4e09\u5c42\u8bc4\u6d4b\u4f53\u7cfb\uff1a\u4ece\u81ea\u52a8\u5316\u6307\u6807\u5230\u7ebf\u4e0a A/B",
             size=14, color=GRAY)

    # ===== Layer 1: Automated metrics =====
    add_text(slide, Inches(0.3), Inches(1.4), Inches(3.7), Inches(0.35),
             "\U0001f4d0 \u7b2c\u4e00\u5c42\uff1a\u81ea\u52a8\u5316\u6307\u6807",
             size=15, bold=True, color=BLUE)

    tbl1 = slide.shapes.add_table(4, 3,
        Inches(0.3), Inches(1.8), Inches(3.7), Inches(1.7))
    t1 = tbl1.table
    for ci, v in enumerate(["\u6307\u6807", "\u542b\u4e49", "\u65b9\u5411"]):
        style_cell(t1.cell(0, ci), v, size=10, bold=True, bg="333399",
                   fg=RGBColor(0xFF, 0xFF, 0xFF))
    rows1 = [
        ["\u56f0\u60d1\u5ea6 (PPL)", "\u6a21\u578b\u5bf9\u9886\u57df\u6587\u672c\u7684\u201c\u719f\u6089\u5ea6\u201d", "\u2193 \u8d8a\u4f4e\u8d8a\u597d"],
        ["\u9ed1\u8bdd\u547d\u4e2d\u7387", "\u751f\u6210\u4e2d\u662f\u5426\u81ea\u7136\u4f7f\u7528\u9886\u57df\u672f\u8bed", "\u2191 \u8d8a\u9ad8\u8d8a\u597d"],
        ["BLEU / ROUGE", "\u4e0e\u53c2\u8003\u7b54\u6848\u7684\u6587\u672c\u91cd\u53e0\u5ea6", "\u2191 \u8d8a\u9ad8\u8d8a\u597d"],
    ]
    for ri, row in enumerate(rows1):
        bg = "F5F5FF" if ri % 2 == 0 else "FFFFFF"
        for ci, v in enumerate(row):
            style_cell(t1.cell(ri + 1, ci), v, size=9, bg=bg)

    # ===== Layer 2: LLM-as-Judge =====
    add_text(slide, Inches(4.2), Inches(1.4), Inches(3.7), Inches(0.35),
             "\U0001f916 \u7b2c\u4e8c\u5c42\uff1aLLM-as-Judge",
             size=15, bold=True, color=BLUE)

    add_multi(slide, Inches(4.2), Inches(1.85), Inches(3.7), Inches(1.7), [
        ("\u7528 GPT-4 / Claude \u5bf9\u8f93\u51fa\u6253\u5206\uff1a", True),
        ("", False),
        ("\u2022 \u6d41\u7545\u5ea6 (1-5\u5206): \u8bed\u53e5\u662f\u5426\u81ea\u7136\u901a\u987a", False),
        ("\u2022 \u9ed1\u8bdd\u81ea\u7136\u5ea6 (1-5\u5206): \u672f\u8bed\u662f\u5426\u7528\u5f97\u5730\u9053", False),
        ("\u2022 \u4e1a\u52a1\u51c6\u786e\u6027 (1-5\u5206): \u5185\u5bb9\u662f\u5426\u7b26\u5408\u7535\u5546\u573a\u666f", False),
        ("", False),
        ("\u4f18\u52bf: \u6210\u672c\u4f4e\u3001\u53ef\u6279\u91cf\u3001\u4e0e\u4eba\u5de5\u76f8\u5173\u6027\u9ad8", True),
    ], size=10, color=GRAY)

    # ===== Layer 3: Human + A/B =====
    add_text(slide, Inches(8.1), Inches(1.4), Inches(3.7), Inches(0.35),
             "\U0001f9d1\u200d\U0001f4bb \u7b2c\u4e09\u5c42\uff1a\u4eba\u5de5 + \u7ebf\u4e0a A/B",
             size=15, bold=True, color=BLUE)

    add_multi(slide, Inches(8.1), Inches(1.85), Inches(3.7), Inches(1.7), [
        ("\u4eba\u5de5\u8bc4\u6d4b:", True),
        ("\u2022 \u672c\u5730\u6bcd\u8bed\u8005\u76f2\u8bc4\uff0c\u91d1\u6807\u51c6\u4f46\u8d35", False),
        ("", False),
        ("\u7ebf\u4e0a A/B \u6d4b\u8bd5:", True),
        ("\u2022 \u90e8\u7f72\u4e24\u7248\u6a21\u578b\uff0c\u770b CTR/\u8f6c\u5316\u7387", False),
        ("\u2022 \u7ec8\u6781\u6307\u6807\uff0c\u4f46\u9700\u8981\u751f\u4ea7\u73af\u5883", False),
    ], size=10, color=GRAY)

    # ===== Demo output example =====
    add_text(slide, Inches(0.3), Inches(3.7), Inches(11.4), Inches(0.35),
             "\U0001f4bb \u672c Demo \u5df2\u5b9e\u73b0\u7684\u8bc4\u6d4b (step5_evaluate.py)",
             size=15, bold=True, color=BLUE)

    # Mock results table
    tbl2 = slide.shapes.add_table(3, 4,
        Inches(0.3), Inches(4.15), Inches(8), Inches(1.0))
    t2 = tbl2.table
    for ci, v in enumerate(["", "Before (\u901a\u7528)", "After (LoRA)", "\u53d8\u5316"]):
        style_cell(t2.cell(0, ci), v, size=11, bold=True, bg="333399",
                   fg=RGBColor(0xFF, 0xFF, 0xFF))
    style_cell(t2.cell(1, 0), "\u56f0\u60d1\u5ea6 (PPL) \u2193", size=11, bg="F5F5FF")
    style_cell(t2.cell(1, 1), "\u2248 15-25", size=11, bg="FFE8E8")
    style_cell(t2.cell(1, 2), "\u2248 3-8", size=11, bg="E8FFE8")
    style_cell(t2.cell(1, 3), "\u2193 60-80%", size=11, bold=True, bg="E8FFE8", fg=GREEN)
    style_cell(t2.cell(2, 0), "\u9ed1\u8bdd\u547d\u4e2d\u7387 \u2191", size=11, bg="FFFFFF")
    style_cell(t2.cell(2, 1), "\u2248 5-15%", size=11, bg="FFE8E8")
    style_cell(t2.cell(2, 2), "\u2248 30-60%", size=11, bg="E8FFE8")
    style_cell(t2.cell(2, 3), "\u2191 +25-45pp", size=11, bold=True, bg="E8FFE8", fg=GREEN)

    # Interpretation
    add_text(slide, Inches(8.5), Inches(4.15), Inches(3.3), Inches(1.0),
             "\U0001f4a1 \u56f0\u60d1\u5ea6\u4e0b\u964d = \u6a21\u578b\u66f4\u201c\u61c2\u201d\u8fd9\u4e2a\u9886\u57df\n"
             "\u9ed1\u8bdd\u547d\u4e2d\u7387\u63d0\u5347 = \u80fd\u81ea\u7136\u8bf4\u201c\u884c\u8bdd\u201d\n\n"
             "\u751f\u4ea7\u73af\u5883\u5efa\u8bae\u4e09\u5c42\u90fd\u505a\uff01",
             size=11, bold=True, color=DARK)

    # ===== Punchline =====
    add_text(slide, Inches(0.3), Inches(5.4), Inches(11.4), Inches(0.9),
             "\U0001f3af \u8bc4\u6d4b\u7b56\u7565\u5efa\u8bae\uff1a"
             "\u5f00\u53d1\u9636\u6bb5\u7528\u81ea\u52a8\u5316\u6307\u6807\u5feb\u901f\u8fed\u4ee3 \u2192 "
             "\u53d1\u5e03\u524d\u7528 LLM-as-Judge \u6279\u91cf\u6253\u5206 \u2192 "
             "\u4e0a\u7ebf\u540e\u7528 A/B \u6d4b\u8bd5\u9a8c\u8bc1\u4e1a\u52a1\u4ef7\u503c\u3002"
             "\u4e09\u5c42\u6f0f\u6597\uff0c\u9010\u6b65\u6536\u7a84\u3002",
             size=13, bold=True, color=DARK)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u8bc4\u6d4b\u662f\u5fae\u8c03\u4e2d\u6700\u5bb9\u6613\u88ab\u5ffd\u7565\u4f46\u6700\u91cd\u8981\u7684\u73af\u8282\u3002"
        "\u5f88\u591a\u56e2\u961f\u5fae\u8c03\u5b8c\u5c31\u4e0a\u7ebf\uff0c\u7f3a\u5c11\u7cfb\u7edf\u6027\u8bc4\u4f30\u3002"
        "\u672c\u6f14\u793a\u5b9e\u73b0\u4e86\u7b2c\u4e00\u5c42\uff08\u56f0\u60d1\u5ea6+\u9ed1\u8bdd\u547d\u4e2d\u7387\uff09\uff0c"
        "\u751f\u4ea7\u73af\u5883\u5efa\u8bae\u4e09\u5c42\u90fd\u505a\u3002"
    )
    return slide


def create_data_roadmap_slide(prs, layout):
    """Slide: data scaling and scenario expansion roadmap."""
    slide = prs.slides.add_slide(layout)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    GREEN = RGBColor(0x22, 0x8B, 0x22)
    ORANGE = RGBColor(0xCC, 0x66, 0x00)

    add_text(slide, Inches(0.5), Inches(0.2), Inches(11), Inches(0.6),
             "\u6570\u636e\u89c4\u6a21\u4e0e\u573a\u666f\u6269\u5c55\u8def\u7ebf\u56fe",
             size=28, bold=True, color=DARK)
    add_text(slide, Inches(0.5), Inches(0.8), Inches(11), Inches(0.4),
             "\u4ece Demo \u73a9\u5177\u5230\u751f\u4ea7\u53ef\u7528\uff0c\u6570\u636e\u662f\u6838\u5fc3\u7684\u201c\u6cb9\u95e8\u201d",
             size=14, color=GRAY)

    # ===== Data volume roadmap table =====
    add_text(slide, Inches(0.3), Inches(1.35), Inches(11), Inches(0.35),
             "\U0001f4c8 \u6570\u636e\u91cf\u8def\u7ebf\u56fe",
             size=16, bold=True, color=BLUE)

    tbl = slide.shapes.add_table(5, 5,
        Inches(0.3), Inches(1.75), Inches(11.4), Inches(2.0))
    t = tbl.table
    headers = ["\u9636\u6bb5", "\u6570\u636e\u91cf",
               "\u9884\u671f\u6548\u679c", "\u8bad\u7ec3\u65f6\u95f4 (1xA100)",
               "\u9002\u7528\u573a\u666f"]
    for ci, v in enumerate(headers):
        style_cell(t.cell(0, ci), v, size=10, bold=True, bg="333399",
                   fg=RGBColor(0xFF, 0xFF, 0xFF))

    rows = [
        ["\U0001f6a9 Demo", "10-30 \u6761",
         "\u8ddf\u901a\u6d41\u7a0b\uff0c\u6548\u679c\u968f\u673a",
         "< 5 \u5206\u949f", "\u6559\u5b66\u6f14\u793a"],
        ["\U0001f680 MVP", "200-500 \u6761",
         "\u57fa\u672c\u98ce\u683c\u8fc1\u79fb",
         "10-30 \u5206\u949f", "\u5185\u90e8\u8bd5\u7528"],
        ["\u2705 \u53ef\u7528\u7ea7", "1000-3000 \u6761",
         "\u7a33\u5b9a\u8f93\u51fa\u9886\u57df\u5185\u5bb9",
         "1-3 \u5c0f\u65f6", "\u5c0f\u8303\u56f4\u4e0a\u7ebf"],
        ["\U0001f31f \u751f\u4ea7\u7ea7", "5000+ \u6761",
         "\u98ce\u683c\u81ea\u7136\u3001\u8986\u76d6\u5168\u9762",
         "3-8 \u5c0f\u65f6", "\u5927\u89c4\u6a21\u90e8\u7f72"],
    ]
    colors = ["FFF5E0", "E8F0FF", "E8FFE8", "F5E8FF"]
    for ri, (row, bg) in enumerate(zip(rows, colors)):
        for ci, v in enumerate(row):
            style_cell(t.cell(ri + 1, ci), v, size=9, bg=bg)

    # ===== Scenario expansion =====
    add_text(slide, Inches(0.3), Inches(3.9), Inches(11), Inches(0.35),
             "\U0001f30f \u573a\u666f\u6269\u5c55\u4e09\u7ef4\u5ea6",
             size=16, bold=True, color=BLUE)

    # Dimension 1: Countries
    add_text(slide, Inches(0.3), Inches(4.35), Inches(3.6), Inches(0.3),
             "\U0001f1ee\U0001f1e9 \u591a\u56fd\u5bb6", size=13, bold=True, color=ORANGE)
    add_multi(slide, Inches(0.3), Inches(4.65), Inches(3.6), Inches(1.3), [
        ("\u5370\u5c3c: cuan, spill, sikat", False),
        ("\u6cf0\u56fd: \u0e25\u0e14\u0e41\u0e23\u0e07, \u0e2a\u0e34\u0e19\u0e04\u0e49\u0e32\u0e14\u0e35", False),
        ("\u8d8a\u5357: deal s\u1ed1c, gi\u1ea3m gi\u00e1 ch\u1ea5n \u0111\u1ed9ng", False),
        ("\u83f2\u5f8b\u5bbe: legit, add to cart na", False),
        ("\u9a6c\u6765: borong, confirm order", False),
    ], size=10, color=GRAY)

    # Dimension 2: Tasks
    add_text(slide, Inches(4.1), Inches(4.35), Inches(3.6), Inches(0.3),
             "\U0001f4dd \u591a\u4efb\u52a1", size=13, bold=True, color=ORANGE)
    add_multi(slide, Inches(4.1), Inches(4.65), Inches(3.6), Inches(1.3), [
        ("\u5546\u54c1\u63a8\u5e7f\u6587\u6848 (\u5f53\u524d\u5df2\u5b9e\u73b0)", True),
        ("\u5ba2\u670d\u8bdd\u672f (\u5f53\u524d\u5df2\u5b9e\u73b0)", True),
        ("\u793e\u4ea4\u5a92\u4f53\u79cd\u8349\u5e16", False),
        ("\u76f4\u64ad\u5e26\u8d27\u811a\u672c", False),
        ("\u63a8\u5ba2\u4f63\u91d1\u7b56\u7565\u5206\u6790", False),
    ], size=10, color=GRAY)

    # Dimension 3: Formats
    add_text(slide, Inches(7.9), Inches(4.35), Inches(3.6), Inches(0.3),
             "\U0001f4f1 \u591a\u683c\u5f0f", size=13, bold=True, color=ORANGE)
    add_multi(slide, Inches(7.9), Inches(4.65), Inches(3.6), Inches(1.3), [
        ("\u77ed\u6587\u6848: TikTok/Reels (\u2264100\u5b57)", False),
        ("\u957f\u6587\u6848: \u535a\u5ba2/\u8be6\u60c5\u9875 (300+\u5b57)", False),
        ("\u5bf9\u8bdd\u4f53: \u5ba2\u670d\u804a\u5929/\u76f4\u64ad\u4e92\u52a8", False),
        ("\u7ed3\u6784\u5316: \u5546\u54c1\u6807\u9898 + \u5356\u70b9\u5217\u8868", False),
        ("\u591a\u8f6e\u5bf9\u8bdd: \u54a8\u8be2\u8ddf\u8fdb\u573a\u666f", False),
    ], size=10, color=GRAY)

    # ===== Punchline =====
    add_text(slide, Inches(0.3), Inches(6.2), Inches(11.4), Inches(0.6),
             "\U0001f4a1 \u5173\u952e\u539f\u5219\uff1a"
             "\u8d28\u91cf > \u6570\u91cf\u3002"
             "100 \u6761\u9ad8\u8d28\u91cf\u6570\u636e > 1000 \u6761\u566a\u58f0\u6570\u636e\u3002"
             "\u6bcf\u4e2a\u56fd\u5bb6/\u4efb\u52a1\u5efa\u8bae\u72ec\u7acb\u8bad\u7ec3\u4e00\u4e2a LoRA \u9002\u914d\u5668\uff0c"
             "\u63a8\u7406\u65f6\u901a\u8fc7 Multi-LoRA \u52a8\u6001\u5207\u6362\u3002",
             size=13, bold=True, color=DARK)

    notes = slide.notes_slide
    notes.notes_text_frame.text = (
        "\u8fd9\u9875\u7ed9\u542c\u4f17\u4e00\u4e2a\u6e05\u6670\u7684\u8def\u7ebf\u56fe\u3002"
        "\u6f14\u793a\u7528\u768430\u6761\u6570\u636e\u53ea\u662f\u8d77\u70b9\uff0c"
        "\u751f\u4ea7\u7ea7\u522b\u9700\u89815000+\u6761\u3002"
        "\u573a\u666f\u6269\u5c55\u4e09\u7ef4\u5ea6\u53ef\u4ee5\u6392\u5217\u7ec4\u5408\uff0c"
        "\u6bcf\u4e2a\u7ec4\u5408\u8bad\u7ec3\u4e00\u4e2a\u72ec\u7acbLoRA\u9002\u914d\u5668\uff0c"
        "\u63a8\u7406\u65f6\u901a\u8fc7Multi-LoRA Batching\u52a8\u6001\u5207\u6362\u3002"
    )
    return slide


def main():
    prs = Presentation("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    layout = find_blank_layout(prs)

    # Find "效果验证" slide index
    target_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and "\u6548\u679c\u9a8c\u8bc1" in shape.text_frame.text:
                target_idx = i
                break
        if target_idx is not None:
            break

    if target_idx is None:
        print("\u274c \u627e\u4e0d\u5230\u201c\u6548\u679c\u9a8c\u8bc1\u201d\u9875\uff0c\u5c06\u63d2\u5165\u5230\u672b\u5c3e")
        target_idx = len(prs.slides) - 1

    print(f"\u627e\u5230\u201c\u6548\u679c\u9a8c\u8bc1\u201d\u9875: Slide {target_idx + 1}")

    sa = create_eval_framework_slide(prs, layout)
    sb = create_data_roadmap_slide(prs, layout)

    # Move after target slide
    slide_list = prs.slides._sldIdLst
    all_els = list(slide_list)
    el_a = all_els[-2]
    el_b = all_els[-1]
    slide_list.remove(el_a)
    slide_list.remove(el_b)

    all_els = list(slide_list)
    ref = all_els[target_idx]
    idx = list(slide_list).index(ref)
    slide_list.insert(idx + 1, el_a)
    slide_list.insert(idx + 2, el_b)

    print(f"\nTotal: {len(prs.slides)} slides\n")
    for i, s in enumerate(prs.slides):
        for shape in s.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:60]
                if t:
                    print(f"  Slide {i+1:2d}: {t}")
                    break

    prs.save("docs/LoRA\u5fae\u8c03\u5b9e\u6218\u6280\u672f\u5206\u4eab_improved.pptx")
    print("\nDone!")


if __name__ == "__main__":
    main()
