"""Add two slides explaining LoRA math more clearly, after Slide 9."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    p.alignment = align
    return tb


def add_multi(slide, left, top, width, height, lines, size=12, color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = is_bold
        if color:
            p.font.color.rgb = color
        p.space_after = Pt(3)
    return tb


def create_dimension_flow_slide(prs, layout):
    """Slide A: visual dimension flow of A and B matrices."""
    sa = prs.slides.add_slide(layout)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)
    RED = RGBColor(0xCC, 0x33, 0x33)
    GREEN = RGBColor(0x22, 0x8B, 0x22)
    ORANGE = RGBColor(0xFF, 0x99, 0x00)

    add_text(sa, Inches(0.5), Inches(0.25), Inches(11), Inches(0.7),
             "\u76f4\u89c9\u7406\u89e3\uff1a\u77e9\u9635 A \u548c B \u5728\u505a\u4ec0\u4e48\uff1f",
             size=28, bold=True, color=DARK)

    add_text(sa, Inches(0.5), Inches(0.95), Inches(11), Inches(0.5),
             "\u6838\u5fc3\u601d\u60f3\uff1a\u9ad8\u7ef4\u7a7a\u95f4\u4e2d\u7684\u4efb\u52a1\u77e5\u8bc6"
             "\u53ea\u9700\u4e00\u4e2a\u4f4e\u7ef4\u5b50\u7a7a\u95f4\u5c31\u80fd\u8868\u8fbe\u3002"
             "A \u8d1f\u8d23\u538b\u7f29\uff0cB \u8d1f\u8d23\u8fd8\u539f\u3002",
             size=13, color=GRAY)

    add_text(sa, Inches(0.5), Inches(1.65), Inches(11), Inches(0.5),
             "\U0001f4d0 \u7ef4\u5ea6\u53d8\u5316\u6d41"
             "\uff08\u4ee5 Qwen-1.8B \u7684\u4e00\u4e2a\u6ce8\u610f\u529b\u5c42\u4e3a\u4f8b"
             "\uff0cd=2048\uff0cr=16\uff09",
             size=16, bold=True, color=BLUE)

    y = Inches(2.2)
    add_text(sa, Inches(0.3), y, Inches(2.0), Inches(1.0),
             "\u8f93\u5165 x\n\u7ef4\u5ea6: [d]\n= [2048]",
             size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(sa, Inches(2.3), Inches(2.5), Inches(1.0), Inches(0.5),
             "\u2192  A  \u2192",
             size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(sa, Inches(3.3), y, Inches(2.2), Inches(1.0),
             "\u77e9\u9635 A (\u964d\u7ef4)\n[r \u00d7 d]\n= [16 \u00d7 2048]",
             size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(sa, Inches(5.5), Inches(2.5), Inches(0.8), Inches(0.5),
             "\u2192",
             size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(sa, Inches(6.0), y, Inches(1.8), Inches(1.0),
             "\u9690\u7a7a\u95f4 h\n\u7ef4\u5ea6: [r]\n= [16]",
             size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(sa, Inches(7.8), Inches(2.5), Inches(1.0), Inches(0.5),
             "\u2192  B  \u2192",
             size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sa, Inches(8.8), y, Inches(2.2), Inches(1.0),
             "\u77e9\u9635 B (\u5347\u7ef4)\n[d \u00d7 r]\n= [2048 \u00d7 16]",
             size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    add_text(sa, Inches(0.5), Inches(3.5), Inches(11), Inches(0.9),
             "\U0001f4a1 \u74f6\u9888\u7ed3\u6784\uff1a"
             "2048 \u2192 16 \u2192 2048\u3002"
             "\u4fe1\u606f\u88ab\u8feb\u901a\u8fc7\u4e00\u4e2a 16 \u7ef4\u7684\u201c\u7a84\u95e8\u201d\uff0c"
             "\u8feb\u4f7f\u6a21\u578b\u53ea\u4fdd\u7559\u6700\u5173\u952e\u7684\u4efb\u52a1\u77e5\u8bc6"
             "\u2014\u2014\u8fd9\u5c31\u662f\u201c\u4f4e\u79e9\u201d\u7684\u529b\u91cf\u3002",
             size=14, bold=True, color=DARK)

    add_text(sa, Inches(0.5), Inches(4.5), Inches(11), Inches(0.4),
             "\U0001f3af \u751f\u6d3b\u7c7b\u6bd4",
             size=16, bold=True, color=BLUE)

    add_multi(sa, Inches(0.5), Inches(4.9), Inches(11), Inches(1.8), [
        ("\u60f3\u8c61\u4f60\u8981\u628a\u4e00\u7bc7 2048 \u5b57\u7684\u6587\u7ae0"
         "\u7ffb\u8bd1\u6210\u53e6\u4e00\u79cd\u8bed\u8a00\uff1a", False),
        ("\u2022 \u5168\u91cf\u5fae\u8c03 = \u9010\u5b57\u9010\u53e5\u91cd\u5199"
         "\uff082048\u00d72048 = 419 \u4e07\u4e2a\u51b3\u7b56\uff09", False),
        ("\u2022 LoRA = \u5148\u538b\u7f29\u6210 16 \u4e2a\u8981\u70b9\uff0c"
         "\u518d\u5c55\u5f00\uff082\u00d72048\u00d716 = 6.5 \u4e07\u4e2a\u51b3\u7b56\uff09", False),
        ("\u5927\u90e8\u5206\u77e5\u8bc6\u53ea\u9700\u6293\u4f4f\u51e0\u4e2a"
         "\u6838\u5fc3\u65b9\u5411\u5c31\u591f\u4e86\uff01", True),
    ], size=12, color=GRAY)

    notes = sa.notes_slide
    notes.notes_text_frame.text = (
        "A\u77e9\u9635\u505a\u964d\u7ef4\uff0c\u628a2048\u7ef4\u538b\u7f29\u523016\u7ef4\u3002"
        "B\u77e9\u9635\u505a\u5347\u7ef4\uff0c\u628a16\u7ef4\u8fd8\u539f\u56de2048\u7ef4\u3002"
        "\u7c7b\u4f3cJPEG\u538b\u7f29\uff1a\u5927\u90e8\u5206\u50cf\u7d20\u662f\u5197\u4f59\u7684\u3002"
    )
    return sa


def create_scaling_slide(prs, layout):
    """Slide B: alpha/r scaling factor explanation."""
    sb = prs.slides.add_slide(layout)
    DARK = RGBColor(0x1A, 0x1A, 0x2E)
    GRAY = RGBColor(0x44, 0x44, 0x44)
    BLUE = RGBColor(0x33, 0x33, 0x99)

    add_text(sb, Inches(0.5), Inches(0.25), Inches(11), Inches(0.7),
             "\u7f29\u653e\u56e0\u5b50\uff1a\u03b1 \u5982\u4f55\u63a7\u5236\u5fae\u8c03\u529b\u5ea6",
             size=28, bold=True, color=DARK)

    add_text(sb, Inches(0.5), Inches(1.05), Inches(11), Inches(0.4),
             "\U0001f4dd \u5b8c\u6574\u516c\u5f0f"
             "\uff08Slide 9 \u7701\u7565\u4e86\u7f29\u653e\u7cfb\u6570\uff09",
             size=16, bold=True, color=BLUE)

    add_text(sb, Inches(1.5), Inches(1.5), Inches(9), Inches(0.7),
             "h  =  W\u2080 \u00b7 x  +  (\u03b1 / r) \u00b7 B \u00b7 A \u00b7 x",
             size=26, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    # Left column: alpha
    add_text(sb, Inches(0.5), Inches(2.5), Inches(5.2), Inches(0.4),
             "\u03b1 (lora_alpha) \u2014 \u7f29\u653e\u7cfb\u6570",
             size=15, bold=True, color=BLUE)
    add_multi(sb, Inches(0.5), Inches(2.9), Inches(5.2), Inches(1.2), [
        ("\u03b1/r \u51b3\u5b9a\u4e86 LoRA \u5c42\u5bf9\u6700\u7ec8\u8f93\u51fa"
         "\u7684\u201c\u97f3\u91cf\u201d\u5927\u5c0f\u3002", False),
        ("\u2022 \u03b1/r \u8d8a\u5927 \u2192 \u5fae\u8c03\u7684\u201c\u58f0\u97f3\u201d"
         "\u8d8a\u54cd\uff0c\u6a21\u578b\u66f4\u591a\u4f9d\u8d56\u65b0\u77e5\u8bc6", False),
        ("\u2022 \u03b1/r \u8d8a\u5c0f \u2192 \u5fae\u8c03\u7684\u201c\u58f0\u97f3\u201d"
         "\u8d8a\u8f7b\uff0c\u6a21\u578b\u66f4\u591a\u4fdd\u7559\u539f\u59cb\u80fd\u529b", False),
    ], size=12, color=GRAY)

    # Right column: rank
    add_text(sb, Inches(6.2), Inches(2.5), Inches(5.2), Inches(0.4),
             "r (rank) \u2014 \u79e9 / \u4fe1\u606f\u5bb9\u91cf",
             size=15, bold=True, color=BLUE)
    add_multi(sb, Inches(6.2), Inches(2.9), Inches(5.2), Inches(1.2), [
        ("r \u51b3\u5b9a\u4e86 LoRA \u80fd\u88c5\u591a\u5c11\u65b0\u77e5\u8bc6\u3002", False),
        ("\u2022 r \u8d8a\u5927 \u2192 \u4f4e\u79e9\u77e9\u9635\u8d8a\u201c\u5bbd\u201d"
         "\uff0c\u5b66\u4e60\u80fd\u529b\u8d8a\u5f3a", False),
        ("\u2022 r \u8d8a\u5c0f \u2192 \u53c2\u6570\u8d8a\u5c11\u8d8a\u9ad8\u6548"
         "\uff0c\u4f46\u53ef\u80fd\u5b66\u4e0d\u5b8c\u6574", False),
    ], size=12, color=GRAY)

    # Config table
    add_text(sb, Inches(0.5), Inches(4.2), Inches(11), Inches(0.4),
             "\u2699\ufe0f \u5e38\u7528\u914d\u7f6e\u53c2\u8003",
             size=16, bold=True, color=BLUE)

    table_shape = sb.shapes.add_table(4, 4,
        Inches(0.5), Inches(4.65), Inches(11), Inches(1.4))
    table = table_shape.table

    data = [
        ["\u914d\u7f6e", "r", "\u03b1", "\u03b1/r (\u5b9e\u9645\u7f29\u653e)"],
        ["\u8f7b\u91cf\u98ce\u683c\u8fc1\u79fb", "8", "16", "2.0"],
        ["\u63a8\u8350\u5e73\u8861\u914d\u7f6e", "16", "32", "2.0"],
        ["\u590d\u6742\u9886\u57df\u77e5\u8bc6", "32", "64", "2.0"],
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = (row_idx == 0)
                p.alignment = PP_ALIGN.CENTER

    add_text(sb, Inches(0.5), Inches(6.15), Inches(11), Inches(0.6),
             "\U0001f4a1 \u7ecf\u9a8c\u6cd5\u5219\uff1a\u4fdd\u6301 \u03b1 = 2r "
             "\u662f\u6700\u5e38\u7528\u7684\u8d77\u70b9\u3002"
             "\u5982\u679c\u6a21\u578b\u51fa\u73b0\u201c\u707e\u96be\u6027\u9057\u5fd8\u201d"
             "\uff0c\u4f18\u5148\u8c03\u5c0f \u03b1 \u800c\u4e0d\u662f r\u3002",
             size=13, bold=True, color=DARK)

    notes = sb.notes_slide
    notes.notes_text_frame.text = (
        "\u8fd9\u9875\u8865\u5145\u4e86Slide 9\u7701\u7565\u7684\u7f29\u653e\u7cfb\u6570\u03b1/r\u3002"
        "\u5982\u679c\u53f0\u4e0b\u6709\u4eba\u95ee\u6a21\u578b\u5b66\u4e86\u65b0\u4e1c\u897f\u4f46\u8bf4\u8bdd"
        "\u53d8\u50bb\u4e86\u600e\u4e48\u529e\uff0c\u7b54\u6848\u5c31\u662f\u8c03\u5c0f\u03b1\u3002"
        "\u7ecf\u5178\u914d\u7f6e\u03b1=2r\u6765\u81ea\u539f\u8bba\u6587\u5b9e\u9a8c\u7ed3\u8bba\u3002"
    )
    return sb


def main():
    prs = Presentation("docs/LoRA微调实战技术分享_improved.pptx")
    layout = find_blank_layout(prs)

    create_dimension_flow_slide(prs, layout)
    create_scaling_slide(prs, layout)

    # Move to after Slide 9 (index 8)
    slide_list = prs.slides._sldIdLst
    all_els = list(slide_list)
    el_a = all_els[-2]
    el_b = all_els[-1]
    slide_list.remove(el_a)
    slide_list.remove(el_b)

    all_els = list(slide_list)
    ref = all_els[8]  # Slide 9
    idx = list(slide_list).index(ref)
    slide_list.insert(idx + 1, el_a)
    slide_list.insert(idx + 2, el_b)

    # Print final order
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()[:55]
                if t:
                    print(f"  Slide {i+1:2d}: {t}")
                    break

    prs.save("docs/LoRA微调实战技术分享_improved.pptx")
    print(f"\nTotal: {len(prs.slides)} slides. Done!")


if __name__ == "__main__":
    main()
