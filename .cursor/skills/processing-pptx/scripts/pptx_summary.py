#!/usr/bin/env python3
"""Print a human-readable summary of a PPTX file."""

import sys

from pptx import Presentation


def summarize(pptx_path: str):
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"File: {pptx_path}")
    print(f"Slides: {total_slides}")
    print(f"Width: {prs.slide_width}, Height: {prs.slide_height}")
    print("=" * 60)

    for i, slide in enumerate(prs.slides, 1):
        print(f"\n--- Slide {i}/{total_slides} ---")
        print(f"Layout: {slide.slide_layout.name}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    preview = text[:200] + "..." if len(text) > 200 else text
                    print(f"  [{shape.name}] {preview}")

            if shape.has_table:
                t = shape.table
                print(f"  [{shape.name}] Table ({len(t.rows)}x{len(t.columns)})")
                if t.rows:
                    header = " | ".join(cell.text for cell in t.rows[0].cells)
                    print(f"    Header: {header[:100]}")

            if shape.shape_type == 13:
                print(f"  [{shape.name}] Image ({shape.image.content_type})")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"  [Notes] {notes[:100]}...")


def main():
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <input.pptx>", file=sys.stderr)
        sys.exit(1)

    summarize(sys.argv[1])


if __name__ == "__main__":
    main()
