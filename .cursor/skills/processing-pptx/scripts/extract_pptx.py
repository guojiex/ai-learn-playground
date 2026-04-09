#!/usr/bin/env python3
"""Extract all text and tables from a PPTX file to JSON."""

import json
import sys
from pathlib import Path

from pptx import Presentation


def extract(pptx_path: str) -> list[dict]:
    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        slide_data: dict = {"slide": i, "shapes": []}

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_data["shapes"].append({"type": "text", "name": shape.name, "text": text})

            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([cell.text for cell in row.cells])
                slide_data["shapes"].append({"type": "table", "name": shape.name, "rows": rows})

            if shape.shape_type == 13:
                slide_data["shapes"].append({
                    "type": "image",
                    "name": shape.name,
                    "content_type": shape.image.content_type,
                    "size": f"{shape.width}x{shape.height}",
                })

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_data["notes"] = notes

        slides.append(slide_data)

    return slides


def main():
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <input.pptx> [output.json]", file=sys.stderr)
        sys.exit(1)

    input_path = sys.argv[1]
    if len(sys.argv) > 2:
        output_path = sys.argv[2]
    else:
        output_path = str(Path(input_path).stem) + "_extracted.json"

    result = extract(input_path)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"Extracted {len(result)} slides to {output_path}")


if __name__ == "__main__":
    main()
